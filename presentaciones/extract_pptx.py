import os
from pptx import Presentation

# Listar archivos pptx ordenados
files = sorted([f for f in os.listdir('.') if f.endswith('.pptx')])

for f in files:
    print(f"\n{'='*70}")
    print(f"=== {f}")
    print(f"{'='*70}")
    prs = Presentation(f)
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                print(shape.text.strip())
