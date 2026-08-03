#!/usr/bin/env python3
"""
Script para mejorar todas las presentaciones de OpenCode con:
- Tema azul pastel predominante
- Mayor variedad visual (no monótono)
- Mejor jerarquía visual
- Elementos decorativos
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import copy

# ==================== PALETA DE COLORES AZUL PASTEL ====================
COLORS = {
    # Azules pastel principales
    'azul_pastel_claro': RGBColor(0xB8, 0xD4, 0xE3),    # #B8D4E3 - Azul claro
    'azul_pastel_medio': RGBColor(0x8E, 0xBD, 0xD8),     # #8EBDD8 - Azul medio
    'azul_pastel_suave': RGBColor(0xA8, 0xD1, 0xDF),     # #A8D1DF - Azul suave
    'azul_cielo': RGBColor(0xCD, 0xE4, 0xF0),            # #CDE4F0 - Cielo
    'azul_ice': RGBColor(0xE8, 0xF1, 0xF8),              # #E8F1F8 - Hielo
    
    # Colores acento
    'naranja_suave': RGBColor(0xF4, 0xA2, 0x61),         # #F4A261 - Naranja suave
    'verde_menta': RGBColor(0x95, 0xD5, 0xB2),           # #95D5B2 - Menta
    'rosa_suave': RGBColor(0xF7, 0xB8, 0xC1),            # #F7B8C1 - Rosa
    'lila_suave': RGBColor(0xD4, 0xB8, 0xE0),            # #D4B8E0 - Lila
    'amarillo_suave': RGBColor(0xFA, 0xE3, 0x9A),        # #FAE39A - Amarillo
    
    # Neutros
    'gris_claro': RGBColor(0xF5, 0xF5, 0xF5),            # #F5F5F5 - Gris claro
    'gris_medio': RGBColor(0x6C, 0x75, 0x7D),            # #6C757D - Gris medio
    'blanco': RGBColor(0xFF, 0xFF, 0xFF),                  # #FFFFFF
    'negro': RGBColor(0x33, 0x33, 0x33),                  # #333333
    'azul_oscuro': RGBColor(0x2C, 0x3E, 0x50),           # #2C3E50 - Azul oscuro
}

# ==================== FUNCIONES DE DISEÑO ====================

def add_gradient_background(slide, color1, color2):
    """Agregar fondo degradado pastel."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1

def add_solid_background(slide, color):
    """Agregar fondo sólido."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_decoration(slide, shape_type, left, top, width, height, fill_color, opacity=0.3):
    """Agregar forma decorativa con opacidad."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle_decoration(slide, left, top, size, color):
    """Agregar círculo decorativo."""
    return add_shape_decoration(
        slide, MSO_SHAPE.OVAL, left, top, size, size, color
    )

def add_accent_bar(slide, left, top, width, height, color):
    """Agregar barra de acento."""
    return add_shape_decoration(
        slide, MSO_SHAPE.RECTANGLE, left, top, width, height, color
    )

def add_bottom_wave(slide, color):
    """Agregar forma decorativa inferior."""
    # Rectángulo decorativo inferior
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.8),
        Inches(10), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_top_accent(slide, color):
    """Agregar barra de acento superior."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_side_accent(slide, color, width=0.3):
    """Agregar barra lateral."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(width), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def style_text_box(textbox, font_name='Calibri', font_size=14, 
                   color=None, bold=False, alignment=PP_ALIGN.LEFT):
    """Estilizar caja de texto."""
    for paragraph in textbox.text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color
            run.font.bold = bold

def create_title_slide(prs, title_text, subtitle_text):
    """Crear diapositiva de título mejorada."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo degradado azul pastel
    add_gradient_background(slide, COLORS['azul_pastel_claro'], COLORS['azul_cielo'])
    
    # Barra superior
    add_top_accent(slide, COLORS['azul_oscuro'])
    
    # Círculos decorativos
    add_circle_decoration(slide, Inches(7.5), Inches(0.5), Inches(1.5), COLORS['naranja_suave'])
    add_circle_decoration(slide, Inches(8.5), Inches(1.5), Inches(0.8), COLORS['verde_menta'])
    add_circle_decoration(slide, Inches(0.5), Inches(5.5), Inches(1.2), COLORS['lila_suave'])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['gris_medio']
    p.alignment = PP_ALIGN.CENTER
    
    # Barra inferior
    add_bottom_wave(slide, COLORS['azul_pastel_suave'])
    
    return slide

def create_content_slide(prs, title_text, content_items, slide_num, accent_color=None):
    """Crear diapositiva de contenido mejorada."""
    if accent_color is None:
        accent_color = COLORS['naranja_suave']
    
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    add_solid_background(slide, COLORS['blanco'])
    
    # Barra lateral azul
    add_side_accent(slide, COLORS['azul_pastel_medio'], 0.25)
    
    # Barra superior con color de acento
    add_accent_bar(slide, Inches(0.25), Inches(0), Inches(9.75), Inches(0.08), accent_color)
    
    # Círculo decorativo
    add_circle_decoration(slide, Inches(8.8), Inches(0.3), Inches(0.8), accent_color)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    
    # Línea separadora
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.1),
        Inches(8.8), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['azul_pastel_suave']
    line.line.fill.background()
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Icono de viñeta
        p.text = f"  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['negro']
        p.space_after = Pt(10)
        p.level = 0
    
    # Pie de página
    add_bottom_wave(slide, COLORS['azul_pastel_claro'])
    
    return slide

def create_two_column_slide(prs, title_text, left_title, left_items, 
                           right_title, right_items, accent_color=None):
    """Crear diapositiva de dos columnas."""
    if accent_color is None:
        accent_color = COLORS['verde_menta']
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    add_solid_background(slide, COLORS['blanco'])
    
    # Barra superior
    add_top_accent(slide, COLORS['azul_pastel_suave'])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    p.alignment = PP_ALIGN.CENTER
    
    # Columna izquierda
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.4),
        Inches(4.5), Inches(5.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['azul_ice']
    left_box.line.color.rgb = COLORS['azul_pastel_medio']
    
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.1), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.1), Inches(4.4))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['negro']
        p.space_after = Pt(8)
    
    # Columna derecha
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.4),
        Inches(4.5), Inches(5.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['blanco']
    right_box.line.color.rgb = accent_color
    
    right_title_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(4.1), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    
    right_content = slide.shapes.add_textbox(Inches(5.4), Inches(2.3), Inches(4.1), Inches(4.4))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['negro']
        p.space_after = Pt(8)
    
    # Pie decorativo
    add_bottom_wave(slide, accent_color)
    
    return slide

def create_bullet_slide(prs, title_text, bullets, accent_color=None):
    """Crear diapositiva con viñetas estilo tarjeta."""
    if accent_color is None:
        accent_color = COLORS['lila_suave']
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    add_solid_background(slide, COLORS['azul_ice'])
    
    # Barra lateral
    add_side_accent(slide, COLORS['azul_pastel_medio'], 0.2)
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    
    # Tarjetas de contenido
    card_colors = [COLORS['blanco'], COLORS['azul_pastel_claro'], COLORS['blanco']]
    
    for i, bullet in enumerate(bullets[:3]):
        y_pos = 1.4 + (i * 1.8)
        
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos),
            Inches(9), Inches(1.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = card_colors[i % 3]
        card.line.color.rgb = COLORS['azul_pastel_medio']
        
        # Icono de acento
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y_pos + 0.3),
            Inches(0.4), Inches(0.4)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = accent_color
        accent_shape.line.fill.background()
        
        # Texto del bullet
        bullet_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.2), Inches(7.8), Inches(1.2))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['negro']
    
    # Pie decorativo
    add_bottom_wave(slide, COLORS['azul_pastel_suave'])
    
    return slide

def create_table_slide(prs, title_text, table_data, accent_color=None):
    """Crear diapositiva con tabla."""
    if accent_color is None:
        accent_color = COLORS['azul_pastel_medio']
    
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    add_solid_background(slide, COLORS['blanco'])
    
    # Barra lateral
    add_side_accent(slide, COLORS['azul_pastel_suave'], 0.2)
    
    # Barra superior
    add_accent_bar(slide, Inches(0.2), Inches(0), Inches(9.8), Inches(0.08), accent_color)
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    
    # Tabla
    if table_data and len(table_data) > 0:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            # Calcular dimensiones
            table_width = min(9.0, cols * 2.0)
            table_height = min(5.5, rows * 0.8)
            
            left = Inches((10 - table_width) / 2)
            top = Inches(1.3)
            
            table_shape = slide.shapes.add_table(
                rows, cols, left, top,
                Inches(table_width), Inches(table_height)
            )
            table = table_shape.table
            
            # Llenar tabla
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        
                        # Estilo de celda
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.font.color.rgb = COLORS['negro']
                        
                        # Color de fondo para header
                        if row_idx == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLORS['azul_pastel_claro']
    
    # Pie decorativo
    add_bottom_wave(slide, COLORS['azul_pastel_suave'])
    
    return slide

def create_closing_slide(prs, title_text, subtitle_text):
    """Crear diapositiva de cierre."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo degradado
    add_gradient_background(slide, COLORS['azul_pastel_medio'], COLORS['azul_pastel_claro'])
    
    # Círculos decorativos
    add_circle_decoration(slide, Inches(1), Inches(1), Inches(2), COLORS['naranja_suave'])
    add_circle_decoration(slide, Inches(7), Inches(5), Inches(1.5), COLORS['verde_menta'])
    add_circle_decoration(slide, Inches(8), Inches(0.5), Inches(1), COLORS['lila_suave'])
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['azul_oscuro']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['gris_medio']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== PROCESAMIENTO DE PRESENTACIONES ====================

def extract_table_data(table):
    """Extraer datos de una tabla."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            text = cell.text.strip()
            cells.append(text)
        rows.append(cells)
    return rows

def extract_content_from_slide(slide):
    """Extraer contenido de una diapositiva existente."""
    title = ""
    content_items = []
    has_table = False
    table_data = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            # El primer texto suele ser el titulo
            if not title and len(text) < 100:
                title = text
            else:
                # Dividir por lineas
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        content_items.append(line)
        
        # Verificar si tiene tabla
        if shape.has_table:
            has_table = True
            table_data = extract_table_data(shape.table)
    
    return title, content_items, has_table, table_data

def improve_presentation(input_path, output_path=None):
    """Mejorar una presentacion existente."""
    if output_path is None:
        output_path = input_path
    
    prs = Presentation(input_path)
    original_slides = list(prs.slides)
    
    # Crear nueva presentacion
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    accent_colors = [
        COLORS['naranja_suave'],
        COLORS['verde_menta'],
        COLORS['lila_suave'],
        COLORS['rosa_suave'],
        COLORS['amarillo_suave']
    ]
    
    for i, slide in enumerate(original_slides):
        title, content_items, has_table, table_data = extract_content_from_slide(slide)
        
        if not title:
            title = f"Diapositiva {i + 1}"
        
        accent = accent_colors[i % len(accent_colors)]
        
        # Determinar tipo de diapositiva
        if i == 0:
            # Primera diapositiva = titulo
            subtitle = content_items[0] if content_items else ""
            create_title_slide(new_prs, title, subtitle)
        elif i == len(original_slides) - 1:
            # Ultima diapositiva = cierre
            create_closing_slide(new_prs, title, "Gracias por su atencion")
        elif has_table:
            # Diapositiva con tabla
            create_table_slide(new_prs, title, table_data, accent)
        elif len(content_items) > 4 and i % 3 == 0:
            # Diapositiva de tarjetas
            create_bullet_slide(new_prs, title, content_items[:3], accent)
        elif i % 4 == 1:
            # Dos columnas (simular)
            half = len(content_items) // 2
            create_two_column_slide(
                new_prs, title,
                "Caracteristicas", content_items[:half] if half > 0 else content_items[:2],
                "Detalles", content_items[half:] if half > 0 else content_items[2:4],
                accent
            )
        else:
            # Contenido regular
            create_content_slide(new_prs, title, content_items, i, accent)
    
    new_prs.save(output_path)
    return output_path

def process_all_presentations():
    """Procesar todas las presentaciones del directorio."""
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    input_dir = os.path.join(base_dir, 'presentaciones')
    output_dir = os.path.join(base_dir, 'presentaciones')
    
    os.makedirs(output_dir, exist_ok=True)
    
    files = sorted([f for f in os.listdir(input_dir) if f.endswith('.pptx')])
    
    print(f"Procesando {len(files)} presentaciones...")
    
    for filename in files:
        input_path = os.path.join(input_dir, filename)
        output_path = os.path.join(output_dir, filename)
        
        try:
            improve_presentation(input_path, output_path)
            print(f"  OK {filename}")
        except Exception as e:
            print(f"  ERROR {filename}: {e}")
    
    print(f"\nProceso completado. Archivos guardados en: {output_dir}")
    """Procesar todas las presentaciones del directorio."""
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    input_dir = os.path.join(base_dir, 'presentaciones')
    output_dir = os.path.join(base_dir, 'presentaciones')
    
    os.makedirs(output_dir, exist_ok=True)
    
    files = sorted([f for f in os.listdir(input_dir) if f.endswith('.pptx')])
    
    print(f"Procesando {len(files)} presentaciones...")
    
    for filename in files:
        input_path = os.path.join(input_dir, filename)
        output_path = os.path.join(output_dir, filename)
        
        try:
            improve_presentation(input_path, output_path)
            print(f"  OK {filename}")
        except Exception as e:
            print(f"  ERROR {filename}: {e}")
    
    print(f"\nProceso completado. Archivos guardados en: {output_dir}")

if __name__ == "__main__":
    process_all_presentations()
