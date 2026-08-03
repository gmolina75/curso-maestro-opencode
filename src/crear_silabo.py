#!/usr/bin/env python3
"""
Presentacion del Silabo Completo - Curso Maestro de OpenCode
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colores de la plantilla
COLOR_TITLE = RGBColor(0x32, 0x57, 0xB8)
COLOR_BODY = RGBColor(0x15, 0x21, 0x3F)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_ACCENT = RGBColor(0x32, 0x57, 0xB8)

FONT_TITLE = 'Roboto Slab'
FONT_BODY = 'Roboto'

# Datos del silabo
MODULES = [
    {"num": 1, "title": "Introduccion y Fundamentos", "topics": ["Que es OpenCode", "Requisitos previos", "Instalacion"]},
    {"num": 2, "title": "Configuracion de Proveedores y Modelos", "topics": ["OpenCode Zen/Go", "Proveedores principales", "Proveedores cloud", "Modelos locales"]},
    {"num": 3, "title": "Configuracion Avanzada", "topics": ["Formato JSON/JSONC", "Ubicaciones de config", "Variables de entorno", "Gestion remota"]},
    {"num": 4, "title": "Interfaz de Usuario TUI", "topics": ["Navegacion basica", "Slash commands", "Atajos de teclado", "Personalizacion"]},
    {"num": 5, "title": "Herramientas Integradas (Tools)", "topics": ["Tools principales", "Configuracion de permisos", "Internals"]},
    {"num": 6, "title": "Modos de Trabajo", "topics": ["Build mode", "Plan mode", "Flujo Plan -> Build"]},
    {"num": 7, "title": "Agentes", "topics": ["Agentes primarios", "Subagentes", "Agentes custom"]},
    {"num": 8, "title": "Agent Skills (SKILL.md)", "topics": ["Ubicacion de skills", "Estructura SKILL.md", "Permisos"]},
    {"num": 9, "title": "AGENTS.md (Contexto)", "topics": ["Que es AGENTS.md", "Contenido tipico", "Uso automatico"]},
    {"num": 10, "title": "MCP Servers", "topics": ["Que es MCP", "MCP local/remoto", "OAuth", "Ejemplos"]},
    {"num": 11, "title": "Custom Tools", "topics": ["Definicion", "Configuracion"]},
    {"num": 12, "title": "Plugins", "topics": ["Que son plugins", "Configuracion", "Ejemplos"]},
    {"num": 13, "title": "Permissions y Policies", "topics": ["Permisos de tools", "Permisos de skills", "Policies experimentales"]},
    {"num": 14, "title": "LSP Servers", "topics": ["Language Server Protocol", "Configuracion", "Herramienta LSP"]},
    {"num": 15, "title": "Formatters", "topics": ["Configuracion", "Formateadores built-in", "Custom formatters"]},
    {"num": 16, "title": "Temas y Personalizacion Visual", "topics": ["Temas de color", "Atajos personalizados"]},
    {"num": 17, "title": "Comandos Custom", "topics": ["Definicion en config", "Archivos markdown"]},
    {"num": 18, "title": "Integracion IDE", "topics": ["VS Code/Cursor/Windsurf", "Atajos IDE", "Context awareness"]},
    {"num": 19, "title": "Integracion GitHub", "topics": ["GitHub Actions", "Eventos soportados", "Uso practico"]},
    {"num": 20, "title": "Integracion GitLab", "topics": ["GitLab Duo", "Self-hosted", "Plugin GitLab"]},
    {"num": 21, "title": "Flujo de Trabajo Plan + Build", "topics": ["Combinacion de modos", "Ejemplos practicos"]},
    {"num": 22, "title": "Agentes Primarios y Subagentes", "topics": ["Arquitectura de agentes", "Navegacion"]},
    {"num": 23, "title": "Agentes Custom Avanzados", "topics": ["Configuracion avanzada", "Casos de uso"]},
    {"num": 24, "title": "Skills Avanzadas", "topics": ["Creacion de skills", "Metadata", "Compatibilidad"]},
    {"num": 25, "title": "AGENTS.md Profundo", "topics": ["Instrucciones adicionales", "Multi-directorio"]},
    {"num": 26, "title": "MCP Avanzado", "topics": ["OAuth avanzado", "Multiples servidores", "Seguridad"]},
    {"num": 27, "title": "Custom Tools y Plugins", "topics": ["Desarrollo de plugins", "Integraciones"]},
    {"num": 28, "title": "Integracion IDE Avanzada", "topics": ["Configuracion detallada", "Multi-IDE"]},
    {"num": 29, "title": "Integracion GitHub Avanzada", "topics": ["Workflows complejos", "Automatizacion"]},
    {"num": 30, "title": "Integracion GitLab Avanzada", "topics": ["Configuracion enterprise", "CI/CD"]},
    {"num": 31, "title": "Troubleshooting y Recursos", "topics": ["Solucion de problemas", "Recursos adicionales", "Comunidad"]}
]

def create_title_slide(prs):
    """Diapositiva de titulo principal."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(8.5), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Silabo Completo"
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.5), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Curso Maestro de OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_BODY
    
    # Descripcion
    desc_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(8.5), Inches(1))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "El Agente de IA de Codigo Abierto para Terminal, Escritorio y IDE"
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_BODY
    
    # Estadisticas
    stats_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.5), Inches(0.5))
    tf = stats_box.text_frame
    p = tf.paragraphs[0]
    p.text = "31 Modulos | 160K+ GitHub Stars | 900+ Contribuidores | 7.5M Desarrolladores"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    return slide

def create_overview_slide(prs):
    """Diapositiva de vista general."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Estructura del Curso"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Columna izquierda - Modulos 1-16
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.2), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    
    for i, mod in enumerate(MODULES[:16]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{mod['num']}. {mod['title']}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(4)
    
    # Columna derecha - Modulos 17-31
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.3), Inches(4.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    
    for i, mod in enumerate(MODULES[16:]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{mod['num']}. {mod['title']}"
        p.font.name = FONT_BODY
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(4)
    
    return slide

def create_module_slide(prs, module):
    """Crear diapositiva para un modulo."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Numero de modulo (grande)
    num_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(1.5), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(module['num'])
    p.font.name = FONT_TITLE
    p.font.size = Pt(48)
    p.font.color.rgb = COLOR_TITLE
    p.font.bold = True
    
    # Titulo del modulo
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.4), Inches(7.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = module['title']
    p.font.name = FONT_TITLE
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TITLE
    
    # Linea separadora
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.3),
        Inches(8.8), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_LIGHT_BG
    line.line.fill.background()
    
    # Temas
    topics_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(4))
    tf = topics_box.text_frame
    tf.word_wrap = True
    
    for i, topic in enumerate(module['topics']):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = topic
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(14)
    
    return slide

def create_section_slide(prs, section_title, module_range):
    """Crear diapositiva de seccion."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Rango
    range_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    tf = range_box.text_frame
    p = tf.paragraphs[0]
    p.text = module_range
    p.font.name = FONT_BODY
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BODY
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_closing_slide(prs):
    """Diapositiva de cierre."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Curso Maestro de OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(40)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "31 Modulos para Dominar OpenCode"
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_BODY
    p.alignment = PP_ALIGN.CENTER
    
    # Recursos
    resources_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = resources_box.text_frame
    p = tf.paragraphs[0]
    p.text = "opencode.ai | github.com/anomalyco/opencode | 160K+ Stars"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    # Crear presentacion
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Diapositiva 1: Titulo
    create_title_slide(prs)
    
    # Diapositiva 2: Vista general
    create_overview_slide(prs)
    
    # Seccion 1: Fundamentos (Modulos 1-6)
    create_section_slide(prs, "Fundamentos", "Modulos 1-6")
    for mod in MODULES[:6]:
        create_module_slide(prs, mod)
    
    # Seccion 2: Agentes y Herramientas (Modulos 7-13)
    create_section_slide(prs, "Agentes y Herramientas", "Modulos 7-13")
    for mod in MODULES[6:13]:
        create_module_slide(prs, mod)
    
    # Seccion 3: Integracion y Configuracion (Modulos 14-20)
    create_section_slide(prs, "Integracion y Configuracion", "Modulos 14-20")
    for mod in MODULES[13:20]:
        create_module_slide(prs, mod)
    
    # Seccion 4: Temas Avanzados (Modulos 21-31)
    create_section_slide(prs, "Temas Avanzados", "Modulos 21-31")
    for mod in MODULES[20:]:
        create_module_slide(prs, mod)
    
    # Diapositiva de cierre
    create_closing_slide(prs)
    
    # Guardar en directorio relativo
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_dir = os.path.join(base_dir, 'presentaciones')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, '00-silabo-completo-opencode.pptx')
    prs.save(output_path)
    print('Silabo creado:', output_path)
    print('Total diapositivas:', len(prs.slides))

if __name__ == '__main__':
    main()
    # Crear presentacion
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Diapositiva 1: Titulo
    create_title_slide(prs)
    
    # Diapositiva 2: Vista general
    create_overview_slide(prs)
    
    # Seccion 1: Fundamentos (Modulos 1-6)
    create_section_slide(prs, "Fundamentos", "Modulos 1-6")
    for mod in MODULES[:6]:
        create_module_slide(prs, mod)
    
    # Seccion 2: Agentes y Herramientas (Modulos 7-13)
    create_section_slide(prs, "Agentes y Herramientas", "Modulos 7-13")
    for mod in MODULES[6:13]:
        create_module_slide(prs, mod)
    
    # Seccion 3: Integracion y Configuracion (Modulos 14-20)
    create_section_slide(prs, "Integracion y Configuracion", "Modulos 14-20")
    for mod in MODULES[13:20]:
        create_module_slide(prs, mod)
    
    # Seccion 4: Temas Avanzados (Modulos 21-31)
    create_section_slide(prs, "Temas Avanzados", "Modulos 21-31")
    for mod in MODULES[20:]:
        create_module_slide(prs, mod)
    
    # Diapositiva de cierre
    create_closing_slide(prs)
    
    # Guardar
    prs.save(output_path)
    print('Silabo creado:', output_path)
    print('Total diapositivas:', len(prs.slides))

if __name__ == '__main__':
    main()
