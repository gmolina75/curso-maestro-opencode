#!/usr/bin/env python3
"""
Script para crear presentaciones de los casos de uso
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colores
COLOR_TITLE = RGBColor(0x32, 0x57, 0xB8)
COLOR_BODY = RGBColor(0x15, 0x21, 0x3F)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

FONT_TITLE = 'Roboto Slab'
FONT_BODY = 'Roboto'

# Casos de uso
USE_CASES = [
    {
        "id": "01",
        "title": "OpenCode + PowerBI",
        "description": "Automatizar reportes PowerBI",
        "slides": [
            {"title": "OpenCode + PowerBI", "content": ["Automatizacion de reportes y dashboards", "Creacion automatica de formulas DAX", "Modelos de datos optimizados"]},
            {"title": "Que es PowerBI + OpenCode", "content": ["PowerBI: Herramienta de business intelligence", "OpenCode: Asistente de IA para codigo", "Combinacion: Automatizacion completa"]},
            {"title": "Casos de Uso", "content": ["Crear modelos de datos", "Generar formulas DAX complejas", "Consultas M para transformaciones", "Reportes automaticos recurrentes"]},
            {"title": "Ejemplo de Prompt", "content": ["@opencode Create a PowerBI data model", "Connect to SQL Server database", "Create DAX measures for:", "- Total Revenue", "- Monthly Growth Rate"]},
            {"title": "Beneficios", "content": ["Ahorro de 60-70% en tiempo", "Formula DAX optimizadas", "Modelos consistentes", "Menos errores humanos"]}
        ]
    },
    {
        "id": "02",
        "title": "OpenCode + GitHub Actions",
        "description": "Pipelines CI/CD automatizados",
        "slides": [
            {"title": "OpenCode + GitHub Actions", "content": ["CI/CD automatizado", "Workflows optimizados", "Seguridad integrada"]},
            {"title": "Que es GitHub Actions", "content": ["Plataforma de CI/CD de GitHub", "Automatizacion de workflows", "Integracion directa con repos"]},
            {"title": "Casos de Uso", "content": ["Crear workflows de testing", "Optimizar pipelines existentes", "Agregar seguridad (SAST, DAST)", "Notificaciones a Slack"]},
            {"title": "Ejemplo", "content": ["@opencode Create a GitHub Actions workflow", "Triggers on push to main", "Runs tests for Python 3.10 and 3.11", "Caches pip dependencies", "Uploads coverage to Codecov"]},
            {"title": "Beneficios", "content": ["Configuracion 5x mas rapida", "Workflows consistentes", "Menos configuracion manual", "Mejores practicas incluidas"]}
        ]
    },
    {
        "id": "03",
        "title": "OpenCode + Docker",
        "description": "Contenedores optimizados",
        "slides": [
            {"title": "OpenCode + Docker", "content": ["Dockerfiles optimizados", "Docker-compose completo", "Seguridad integrada"]},
            {"title": "Que es Docker", "content": ["Plataforma de contenedores", "Empaquetar aplicaciones", "Despliegue consistente"]},
            {"title": "Casos de Uso", "content": ["Crear Dockerfiles multi-stage", "Docker-compose con servicios", "Optimizar imagenes existentes", ".dockerignore apropiado"]},
            {"title": "Ejemplo", "content": ["@opencode Create multi-stage Dockerfile", "Node.js 20 backend", "Multi-stage for smaller image", "Non-root user", "Health checks included"]},
            {"title": "Beneficios", "content": ["Imagenes 40-60% mas pequenas", "Mejores practicas de seguridad", "Configuracion consistente", "Menos debugging"]}
        ]
    },
    {
        "id": "04",
        "title": "OpenCode + React",
        "description": "Componentes frontend modernos",
        "slides": [
            {"title": "OpenCode + React", "content": ["Componentes reutilizables", "Hooks personalizados", "TypeScript completo"]},
            {"title": "Que es React", "content": ["Biblioteca de UI de Facebook", "Componentes basados en estado", "Ecosistema amplio"]},
            {"title": "Casos de Uso", "content": ["Crear componentes completos", "Hooks personalizados", "Refactorizar codigo", "Tests y Storybook"]},
            {"title": "Ejemplo", "content": ["@opencode Create DataTable component", "- Sorting by columns", "- Pagination", "- Search/filter", "- Export to CSV"]},
            {"title": "Beneficios", "content": ["Componentes mas reutilizables", "Codigo mas mantenible", "Tests automaticos", "Documentacion incluida"]}
        ]
    },
    {
        "id": "05",
        "title": "OpenCode + Data Science",
        "description": "Analisis de datos y ML",
        "slides": [
            {"title": "OpenCode + Data Science", "content": ["Scripts de analisis", "Modelos de ML", "Visualizaciones"]},
            {"title": "Que es Data Science", "content": ["Analisis de datos con Python", "Machine Learning", "Pandas, NumPy, Scikit-learn"]},
            {"title": "Casos de Uso", "content": ["Analisis exploratorio (EDA)", "Modelos de ML", "Limpieza de datos", "Visualizaciones"]},
            {"title": "Ejemplo", "content": ["@opencode Create EDA script", "Load CSV with pandas", "Handle missing values", "Generate statistical summary", "Create visualizations"]},
            {"title": "Beneficios", "content": ["Scripts mas rapidos", "Mejores practicas de ML", "Codigo documentado", "Reproducibilidad"]}
        ]
    },
    {
        "id": "06",
        "title": "OpenCode + AWS",
        "description": "Infraestructura cloud",
        "slides": [
            {"title": "OpenCode + AWS", "content": ["Lambda functions", "Terraform", "CI/CD en AWS"]},
            {"title": "Que es AWS", "content": ["Amazon Web Services", "Cloud computing", "Servicios escalables"]},
            {"title": "Casos de Uso", "content": ["Crear Lambda functions", "Terraform para infraestructura", "CI/CD para AWS", "Seguridad IAM"]},
            {"title": "Ejemplo", "content": ["@opencode Create Lambda function", "API Gateway trigger", "DynamoDB integration", "IAM minimal permissions", "Error handling"]},
            {"title": "Beneficios", "content": ["Configuracion mas rapida", "Mejores practicas de seguridad", "Costos optimizados", "Infraestructura como codigo"]}
        ]
    },
    {
        "id": "07",
        "title": "OpenCode + GitLab",
        "description": "DevOps completo con GitLab",
        "slides": [
            {"title": "OpenCode + GitLab", "content": ["GitLab CI/CD", "Auto DevOps", "Integracion completa"]},
            {"title": "Que es GitLab", "content": ["Plataforma DevOps", "Repositorio + CI/CD", "Gestion de proyecto"]},
            {"title": "Casos de Uso", "content": ["Pipelines completos", "Auto DevOps", "Integracion con GitLab", "Container registry"]},
            {"title": "Ejemplo", "content": ["@opencode Create .gitlab-ci.yml", "Build stage (Docker)", "Test stage (parallel)", "Security scanning", "Deploy to staging"]},
            {"title": "Beneficios", "content": ["Pipelines mas rapidos", "Menos configuracion manual", "Mejores practicas CI/CD", "Integracion completa"]}
        ]
    },
    {
        "id": "08",
        "title": "OpenCode + Bases de Datos",
        "description": "Esquemas y queries SQL",
        "slides": [
            {"title": "OpenCode + Bases de Datos", "content": ["Esquemas optimizados", "Migraciones", "Queries performantes"]},
            {"title": "Que es SQL", "content": ["Lenguaje de consultas", "PostgreSQL, MySQL, SQLite", "Relational databases"]},
            {"title": "Casos de Uso", "content": ["Disenar esquemas", "Crear migraciones", "Optimizar queries", "Documentar schema"]},
            {"title": "Ejemplo", "content": ["@opencode Design PostgreSQL schema", "E-commerce platform", "Users with roles", "Products with inventory", "Orders with payments"]},
            {"title": "Beneficios", "content": ["Esquemas mejor diseniados", "Queries mas rapidas", "Migraciones sin errores", "Documentacion automatica"]}
        ]
    },
    {
        "id": "09",
        "title": "OpenCode + API REST",
        "description": "APIs completas con docs",
        "slides": [
            {"title": "OpenCode + API REST", "content": ["APIs completas", "Documentacion OpenAPI", "Tests automaticos"]},
            {"title": "Que es una API REST", "content": ["Interfaz de programacion", "Endpoints HTTP", "JSON como formato"]},
            {"title": "Casos de Uso", "content": ["Crear APIs completas", "Autenticacion JWT", "Documentacion Swagger", "Tests de integracion"]},
            {"title": "Ejemplo", "content": ["@opencode Create REST API", "Express.js + TypeScript", "JWT authentication", "Input validation", "Swagger documentation"]},
            {"title": "Beneficios", "content": ["API mas rapida de crear", "Documentacion actualizada", "Tests comprehensivos", "Seguridad integrada"]}
        ]
    },
    {
        "id": "10",
        "title": "OpenCode + Automatizacion",
        "description": "Scripts y tareas automaticas",
        "slides": [
            {"title": "OpenCode + Automatizacion", "content": ["Scripts de automatizacion", "Scrapers web", "Monitoreo de servicios"]},
            {"title": "Que es Automatizacion", "content": ["Tareas repetitivas automatizadas", "Scripts programados", "Ahorro de tiempo"]},
            {"title": "Casos de Uso", "content": ["Web scrapers", "Generacion de reportes", "Monitoreo de APIs", "Notificaciones automaticas"]},
            {"title": "Ejemplo", "content": ["@opencode Create web scraper", "Extract product prices", "Handle pagination", "Store in CSV", "Runs daily via cron"]},
            {"title": "Beneficios", "content": ["Automatizacion mas rapida", "Scripts robustos", "Manejo de errores", "Documentacion completa"]}
        ]
    }
]

def create_use_case_presentation(use_case, output_dir):
    """Crear presentacion para un caso de uso."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for i, slide_data in enumerate(use_case['slides']):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Fondo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_WHITE
        
        # Titulo
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.5), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data['title']
        p.font.name = FONT_TITLE
        p.font.size = Pt(28)
        p.font.color.rgb = COLOR_TITLE
        
        # Contenido
        content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.5), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for j, item in enumerate(slide_data['content']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_BODY
            p.space_after = Pt(10)
    
    # Guardar
    output_path = os.path.join(output_dir, f"{use_case['id']}-{use_case['title'].lower().replace(' ', '-').replace('+', '')}.pptx")
    prs.save(output_path)
    return output_path

def main():
    # Directorio de salida relativo al script
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_dir = os.path.join(base_dir, 'presentaciones', 'casos-de-uso')
    os.makedirs(output_dir, exist_ok=True)
    
    print(f"Creando {len(USE_CASES)} presentaciones de casos de uso...")
    print()
    
    for use_case in USE_CASES:
        try:
            path = create_use_case_presentation(use_case, output_dir)
            print(f"OK: {use_case['title']}")
        except Exception as e:
            print(f"ERROR: {use_case['title']}: {e}")
    
    print()
    print("Proceso completado!")
    # Directorio de salida
    output_dir = r"G:\My Drive\estudios\cursos\OpenCode\presentaciones_mejoradas\casos-de-uso"
    os.makedirs(output_dir, exist_ok=True)
    
    print(f"Creando {len(USE_CASES)} presentaciones de casos de uso...")
    print()
    
    for use_case in USE_CASES:
        try:
            path = create_use_case_presentation(use_case, output_dir)
            print(f"OK: {use_case['title']}")
        except Exception as e:
            print(f"ERROR: {use_case['title']}: {e}")
    
    print()
    print("Proceso completado!")

if __name__ == "__main__":
    main()
