#!/usr/bin/env python3
"""
Crear presentacion 01 usando el estilo de la plantilla
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colores de la plantilla
COLOR_TITLE = RGBColor(0x32, 0x57, 0xB8)    # #3257B8 - Azul
COLOR_BODY = RGBColor(0x15, 0x21, 0x3F)     # #15213F - Navy oscuro
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA) # Fondo claro
COLOR_ACCENT = RGBColor(0x32, 0x57, 0xB8)   # Mismo azul para acentos

# Fuentes de la plantilla
FONT_TITLE = 'Roboto Slab'
FONT_BODY = 'Roboto'

def create_slide_01(prs):
    """Diapositiva 1: Titulo principal"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(5.5), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Introduccion a OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(38)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(5.5), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Modulo 1 | Duracion: 45 minutos"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_BODY
    p.alignment = PP_ALIGN.LEFT
    
    # Autor
    author_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(5.5), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Curso OpenCode"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_BODY
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def create_slide_02(prs):
    """Diapositiva 2: Que es OpenCode"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo de seccion (izquierda arriba)
    section_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(4), Inches(0.6))
    tf = section_box.text_frame
    p = tf.paragraphs[0]
    p.text = "QUE ES OPENCODE"
    p.font.name = FONT_TITLE
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TITLE
    p.font.bold = True
    
    # Titulo principal
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(5.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Agente de IA de Codigo Abierto"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Contenido
    content_items = [
        "Terminal First: Ejecuta comandos directamente en tu terminal",
        "Multi-Proveedor: Soporta mas de 75 proveedores de IA",
        "Codigo Abierto: Licenciado bajo MIT",
        "Privacidad First: Tus datos se procesan con el proveedor que eliges"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(5.5), Inches(3))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(12)
    
    return slide

def create_slide_03(prs):
    """Diapositiva 3: Arquitectura"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Arquitectura de OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Componentes
    components = [
        ("TUI (Terminal UI)", "Interfaz interactiva en terminal para desarrollo diario"),
        ("Desktop App", "Aplicacion de escritorio nativa para usuarios GUI"),
        ("IDE Extension", "Extensiones para VS Code y JetBrains"),
        ("CLI Mode", "Modo de lineas de comandos para automatizacion")
    ]
    
    for i, (name, desc) in enumerate(components):
        y_pos = 1.4 + (i * 1.2)
        
        # Cuadro de componente
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6), Inches(y_pos),
            Inches(8.5), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_LIGHT_BG
        box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        
        # Nombre
        name_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.1), Inches(3), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE
        
        # Descripcion
        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.5), Inches(8), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_BODY
    
    return slide

def create_slide_04(prs):
    """Diapositiva 4: Comparativa"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "OpenCode vs Competidores"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Tabla de comparativa
    rows = 6
    cols = 4
    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.6), Inches(1.3),
        Inches(8.5), Inches(5)
    )
    table = table_shape.table
    
    # Headers
    headers = ["Caracteristica", "OpenCode", "Claude Code", "Cursor"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            para.font.name = FONT_BODY
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = COLOR_WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
    
    # Datos
    data = [
        ["Codigo Abierto", "Si (MIT)", "No", "No"],
        ["Multi-Proveedor", "Si (75+)", "Solo Anthropic", "No"],
        ["Terminal First", "Si", "Si", "No"],
        ["Opciones Locales", "Si (Ollama)", "No", "No"],
        ["Gratuito", "Si (pago por uso)", "No", "No"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                para.font.name = FONT_BODY
                para.font.size = Pt(11)
                para.font.color.rgb = COLOR_BODY
            
            # Color de fondo alternado
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_BG
    
    return slide

def create_slide_05(prs):
    """Diapositiva 5: Ventajas"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ventajas Clave de OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Ventajas
    advantages = [
        ("1. Modelo Open Source", "Codigo auditable, sin vendor lock-in, comunidad activa"),
        ("2. Multi-Proveedor", "Cambia de modelo sin cambiar de herramienta"),
        ("3. Privacidad y Cumplimiento", "Datos directos al proveedor, opciones locales"),
        ("4. Rendimiento y Eficiencia", "TUI optimizada, streaming en tiempo real")
    ]
    
    for i, (title, desc) in enumerate(advantages):
        y_pos = 1.3 + (i * 1.3)
        
        # Titulo de ventaja
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(8), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE
        
        # Descripcion
        desc_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.4), Inches(8), Inches(0.6))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_BODY
    
    return slide

def create_slide_06(prs):
    """Diapositiva 6: Demo"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Demo: Explorando OpenCode"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Codigo de ejemplo
    code_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.5), Inches(4.5))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    code_lines = [
        "# La experiencia basica de OpenCode",
        "opencode",
        "",
        "# Esto abre la interfaz TUI donde puedes:",
        "# - Escribir prompts en lenguaje natural",
        "# - Ver el codigo que la IA genera en tiempo real",
        "# - Aprobar o rechazar cambios",
        "# - Ejecutar comandos de terminal"
    ]
    
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = 'Consolas'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_BODY
    
    return slide

def create_slide_07(prs):
    """Diapositiva 7: Resumen"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Resumen"
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Puntos clave
    key_points = [
        "OpenCode es el agente de IA open source mas popular",
        "Multi-proveedor con soporte para 75+ proveedores",
        "Privacidad y cumplimiento normativo",
        "Terminal First con TUI optimizada"
    ]
    
    points_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(3))
    tf = points_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(key_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.name = FONT_BODY
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BODY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)
    
    return slide

def main():
    # Crear presentacion
    prs = Presentation()
    
    # Tamano de diapositiva (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Crear todas las diapositivas
    create_slide_01(prs)
    create_slide_02(prs)
    create_slide_03(prs)
    create_slide_04(prs)
    create_slide_05(prs)
    create_slide_06(prs)
    create_slide_07(prs)
    
    # Guardar
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_dir = os.path.join(base_dir, 'presentaciones')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, '01-introduccion-opencode.pptx')
    prs.save(output_path)
    print('Presentacion creada:', output_path)

if __name__ == '__main__':
    main()
