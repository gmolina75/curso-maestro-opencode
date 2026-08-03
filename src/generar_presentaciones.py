#!/usr/bin/env python3
"""
Script para generar 31 presentaciones PowerPoint del Curso Maestro de OpenCode
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colores del tema OpenCode
COLORS = {
    'primary': RGBColor(0x21, 0x1E, 0x1E),      # Negro oscuro
    'secondary': RGBColor(0x65, 0x63, 0x63),     # Gris medio
    'accent': RGBColor(0xCF, 0xCE, 0xCD),        # Gris claro
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF1, 0xEC, 0xEC),
    'code_bg': RGBColor(0x2D, 0x2D, 0x2D),
    'green': RGBColor(0x4E, 0xC9, 0xB0),
    'blue': RGBColor(0x56, 0x9C, 0xD6),
    'orange': RGBColor(0xCE, 0x91, 0x78),
}

# Contenido de cada presentación
PRESENTATIONS = [
    {
        "file": "01-introduccion-opencode.pptx",
        "title": "Introducción a OpenCode",
        "module": 1,
        "duration": "1.5 horas",
        "objectives": [
            "Comprender qué es OpenCode y su arquitectura",
            "Conocer las ventajas sobre otros agentes de IA",
            "Identificar casos de uso apropiados",
            "Explorar el ecosistema de OpenCode"
        ],
        "slides": [
            {"type": "title", "content": "Clase 1: Introducción a OpenCode"},
            {"type": "content", "title": "¿Qué es OpenCode?", "bullets": [
                "Agente de IA de código abierto para coding",
                "Disponible como: Terminal TUI, Desktop App, IDE Extension",
                "Construido en Go - rápido y eficiente",
                "160K+ GitHub Stars, 900+ contribuidores",
                "7.5M desarrolladores activos mensuales"
            ]},
            {"type": "content", "title": "Arquitectura de OpenCode", "bullets": [
                "Terminal User Interface (TUI) - interfaz interactiva",
                "Desktop App - aplicación de escritorio (beta)",
                "IDE Extension - VS Code, Cursor, Windsurf",
                "CLI - modo headless para automatización",
                "Server - modo servidor para integración"
            ]},
            {"type": "comparison", "title": "OpenCode vs Competidores", "headers": ["Característica", "OpenCode", "Claude Code", "Cursor"],
             "rows": [
                 ["Open Source", "✅ Sí", "❌ No", "❌ No"],
                 ["Precio", "🆓 Gratis", "💰 Pagado", "💰 Pagado"],
                 ["Multi-provider", "✅ 75+", "❌ Solo Anthropic", "❌ Limitado"],
                 ["Self-hosted", "✅ Sí", "❌ No", "❌ No"],
                 ["Privacy", "✅ Local", "⚠️ Servidor", "⚠️ Servidor"]
             ]},
            {"type": "content", "title": "Ventajas Clave", "bullets": [
                "Libre y open source - sin vendor lock-in",
                "Multi-provider: Claude, GPT, Gemini, modelos locales",
                "Privacy-first: código nunca sale de tu máquina",
                "Altamente customizable: agents, skills, plugins",
                "LSP integrado: inteligencia de código nativa",
                "Multi-sesión: trabajar en paralelo"
            ]},
            {"type": "demo", "title": "Demo: Explorando OpenCode", "steps": [
                "Visitar opencode.ai",
                "Explorar la documentación",
                "Ver el repositorio en GitHub",
                "Unir al Discord community"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "OpenCode es el agente de IA open source más popular",
                "Multi-plataforma: terminal, desktop, IDE",
                "Privacy-first con opciones de modelos locales",
                "Comunidad activa y en crecimiento"
            ]}
        ]
    },
    {
        "file": "02-instalacion.pptx",
        "title": "Instalación de OpenCode",
        "module": 2,
        "duration": "1 hora",
        "objectives": [
            "Instalar OpenCode en cualquier plataforma",
            "Verificar la instalación correctamente",
            "Actualizar OpenCode cuando sea necesario",
            "Resolver problemas comunes de instalación"
        ],
        "slides": [
            {"type": "title", "content": "Clase 2: Instalación de OpenCode"},
            {"type": "content", "title": "Prerrequisitos", "bullets": [
                "Terminal moderno: WezTerm, Alacritty, Ghostty, Kitty",
                "API keys de proveedores LLM (o usar OpenCode Zen)",
                "Conocimientos básicos de línea de comandos",
                "Git (recomendado)"
            ]},
            {"type": "code", "title": "Instalación Recomendada", "code": "curl -fsSL https://opencode.ai/install | bash"},
            {"type": "code", "title": "Instalación con Node.js", "code": "# npm\nnpm install -g opencode-ai\n\n# bun\nbun install -g opencode-ai\n\n# pnpm\npnpm install -g opencode-ai\n\n# yarn\nyarn global add opencode-ai"},
            {"type": "code", "title": "Instalación con Homebrew", "code": "brew install anomalyco/tap/opencode"},
            {"type": "content", "title": "Instalación en Windows", "bullets": [
                "Chocolatey: choco install opencode",
                "Scoop: scoop install opencode",
                "Mise: mise use -g github:anomalyco/opencode",
                "Docker: docker run -it --rm ghcr.io/anomalyco/opencode",
                "WSL (recomendado para mejor experiencia)"
            ]},
            {"type": "code", "title": "Verificación", "code": "opencode --version"},
            {"type": "summary", "title": "Resumen", "points": [
                "Múltiples métodos de instalación disponibles",
                "WSL recomendado para Windows",
                "Verificar siempre después de instalar"
            ]}
        ]
    },
    {
        "file": "03-opencode-zen-go.pptx",
        "title": "OpenCode Zen y Go",
        "module": 3,
        "duration": "1 hora",
        "objectives": [
            "Configurar OpenCode Zen para empezar rápido",
            "Entender el modelo de pago pay-as-you-go",
            "Comparar Zen vs Go vs BYOK",
            "Seleccionar el plan adecuado"
        ],
        "slides": [
            {"type": "title", "content": "Clase 3: OpenCode Zen y Go"},
            {"type": "content", "title": "¿Qué es OpenCode Zen?", "bullets": [
                "Lista curada de modelos verificados por OpenCode",
                "Pay-as-you-go: $20 saldo, auto top-up",
                "Sin markup en inferencia",
                "Modelos probados y optimizados para coding",
                "Sin necesidad de probar 40+ combinaciones"
            ]},
            {"type": "code", "title": "Configuración de Zen", "code": "# 1. Ejecutar /connect en TUI\n/connect\n\n# 2. Seleccionar 'OpenCode Zen'\n# 3. Ir a opencode.ai/auth\n# 4. Iniciar sesión y agregar datos de facturación\n# 5. Copiar API key\n# 6. Pegar la API key en el prompt"},
            {"type": "content", "title": "OpenCode Go", "bullets": [
                "Plan de suscripción de bajo costo",
                "Modelos de código abierto populares",
                "Acceso confiable a modelos probados",
                "Misma configuración que Zen"
            ]},
            {"type": "comparison", "title": "Comparativa de Planes", "headers": ["Característica", "Zen", "Go", "BYOK"],
             "rows": [
                 ["Modelos", "Curados", "Open source", "Cualquiera"],
                 ["Pago", "Pay-as-you-go", "Suscripción", "API keys propias"],
                 ["Configuración", "Fácil", "Fácil", "Manual"],
                 ["Soporte", "OpenCode team", "OpenCode team", "Propio"]
             ]},
            {"type": "demo", "title": "Demo: Configurar Zen", "steps": [
                "Ejecutar opencode en terminal",
                "Ejecutar /connect",
                "Seleccionar OpenCode Zen",
                "Seguir instrucciones en navegador",
                "Pegar API key",
                "Ejecutar /models para ver modelos"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Zen es la mejor opción para empezar",
                "Go es ideal para modelos open source",
                "BYOK para usuarios avanzados"
            ]}
        ]
    },
    {
        "file": "04-proveedores-principales.pptx",
        "title": "Proveedores Principales",
        "module": 4,
        "duration": "1.5 horas",
        "objectives": [
            "Configurar Anthropic (Claude) con OAuth",
            "Configurar OpenAI (GPT) con OAuth",
            "Configurar GitHub Copilot y GitLab Duo",
            "Configurar Google Vertex AI"
        ],
        "slides": [
            {"type": "title", "content": "Clase 4: Proveedores Principales"},
            {"type": "content", "title": "Anthropic (Claude)", "bullets": [
                "Soporte para Claude Pro/Max (OAuth)",
                "Modelos: Claude Opus, Sonnet, Haiku",
                "Configuración con /connect",
                "Acceso directo a suscripción existente"
            ]},
            {"type": "code", "title": "Configurar Anthropic", "code": "/connect\n# Seleccionar Anthropic\n# Seleccionar 'Claude Pro/Max'\n# Autenticar en navegador\n# ¡Listo! Modelos disponibles en /models"},
            {"type": "content", "title": "OpenAI (GPT)", "bullets": [
                "Soporte para ChatGPT Plus/Pro (OAuth)",
                "Modelos: GPT-4o, GPT-5, o1, o3",
                "Misma configuración que Anthropic",
                "Usa suscripción existente"
            ]},
            {"type": "content", "title": "GitHub Copilot", "bullets": [
                "Autenticación OAuth",
                "Navegar a github.com/login/device",
                "Ingresar código de verificación",
                "Acceso a modelos de GitHub"
            ]},
            {"type": "content", "title": "GitLab Duo", "bullets": [
                "Requiere suscripción Premium o Ultimate",
                "OAuth o Personal Access Token",
                "Modelos: haiku-4-5, sonnet-4-5, opus-4-5",
                "Soporte para self-hosted GitLab"
            ]},
            {"type": "content", "title": "Google Vertex AI", "bullets": [
                "Requiere Google Cloud project",
                "Modelos Gemini disponibles",
                "Configuración con variables de entorno",
                "Autenticación con service account"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Anthropic y OpenAI usan OAuth (fácil)",
                "GitHub Copilot y GitLab Duo también soportados",
                "Google Vertex AI para modelos Gemini",
                "Todos se configuran con /connect"
            ]}
        ]
    },
    {
        "file": "05-proveedores-cloud.pptx",
        "title": "Proveedores Cloud",
        "module": 5,
        "duration": "1.5 horas",
        "objectives": [
            "Configurar proveedores cloud principales",
            "Entender las diferencias entre proveedores",
            "Usar gateways de IA como Cloudflare",
            "Optimizar costos con routing inteligente"
        ],
        "slides": [
            {"type": "title", "content": "Clase 5: Proveedores Cloud"},
            {"type": "content", "title": "Amazon Bedrock (AWS)", "bullets": [
                "Acceso a múltiples modelos (Claude, Llama, etc.)",
                "Configuración con AWS credentials",
                "Soporte para VPC endpoints",
                "Región y profile configurables"
            ]},
            {"type": "code", "title": "Configurar Bedrock", "code": "# Variables de entorno\nAWS_PROFILE=my-profile AWS_REGION=us-east-1 opencode\n\n# O en opencode.json\n{\n  \"provider\": {\n    \"amazon-bedrock\": {\n      \"options\": {\n        \"region\": \"us-east-1\",\n        \"profile\": \"my-aws-profile\"\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "Azure OpenAI", "bullets": [
                "Recursos de Azure OpenAI",
                "Deploy models en Azure AI Foundry",
                "Variable AZURE_RESOURCE_NAME",
                "Content filter: usar Default no DefaultV2"
            ]},
            {"type": "content", "title": "Cloudflare AI Gateway", "bullets": [
                "Endpoint unificado para múltiples proveedores",
                "Unified Billing - sin API keys separadas",
                "Account ID y Gateway ID requeridos",
                "Soporte para OpenAI, Anthropic, Workers AI"
            ]},
            {"type": "content", "title": "Otros Proveedores", "bullets": [
                "DigitalOcean: Inference Routers para routing inteligente",
                "Fireworks AI: Modelos open source optimizados",
                "Groq: Ultra rápido para inferencia",
                "Hugging Face: 17+ providers",
                "NVIDIA: Modelos Nemotron",
                "OpenRouter: Acceso a múltiples proveedores"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "AWS Bedrock para enterprise",
                "Azure para organizaciones Microsoft",
                "Cloudflare para simplificar billing",
                "DigitalOcean para routing inteligente"
            ]}
        ]
    },
    {
        "file": "06-modelos-locales.pptx",
        "title": "Modelos Locales",
        "module": 6,
        "duration": "1 hora",
        "objectives": [
            "Configurar Ollama con OpenCode",
            "Usar LM Studio y llama.cpp",
            "Optimizar tool calls con modelos locales",
            "Entender limitaciones y mejores prácticas"
        ],
        "slides": [
            {"type": "title", "content": "Clase 6: Modelos Locales"},
            {"type": "content", "title": "¿Por qué Modelos Locales?", "bullets": [
                "Privacidad total - código nunca sale de tu máquina",
                "Sin costos de inferencia",
                "Funciona sin internet",
                "Control total sobre el modelo"
            ]},
            {"type": "code", "title": "Configurar Ollama", "code": "# Instalar Ollama\n# https://ollama.com\n\n# Configurar en opencode.json\n{\n  \"provider\": {\n    \"ollama\": {\n      \"npm\": \"@ai-sdk/openai-compatible\",\n      \"name\": \"Ollama (local)\",\n      \"options\": {\n        \"baseURL\": \"http://localhost:11434/v1\"\n      },\n      \"models\": {\n        \"llama2\": { \"name\": \"Llama 2\" }\n      }\n    }\n  }\n}"},
            {"type": "code", "title": "Configurar LM Studio", "code": "# LM Studio ejecuta en puerto 1234\n{\n  \"provider\": {\n    \"lmstudio\": {\n      \"npm\": \"@ai-sdk/openai-compatible\",\n      \"name\": \"LM Studio (local)\",\n      \"options\": {\n        \"baseURL\": \"http://127.0.0.1:1234/v1\"\n      },\n      \"models\": {\n        \"google/gemma-3n-e4b\": {\n          \"name\": \"Gemma 3n-e4b (local)\"\n        }\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "Tips para Tool Calls", "bullets": [
                "Aumentar num_ctx en Ollama (16k-32k)",
                "Usar modelos con soporte de tool calling",
                "Qwen-Coder y DeepSeek-Coder son buenos",
                "Verificar que el servidor esté corriendo"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Ollama es la opción más popular",
                "LM Studio para GUI amigable",
                "llama.cpp para máximo control",
                "Aumentar contexto para mejor tool calling"
            ]}
        ]
    },
    {
        "file": "07-proveedor-custom.pptx",
        "title": "Proveedor Custom y Gestión de Modelos",
        "module": 7,
        "duration": "1 hora",
        "objectives": [
            "Crear proveedores personalizados",
            "Configurar baseURL para proxies",
            "Gestionar modelos con blacklist/whitelist",
            "Optimizar timeouts y rendimiento"
        ],
        "slides": [
            {"type": "title", "content": "Clase 7: Proveedor Custom"},
            {"type": "code", "title": "Estructura de Proveedor Custom", "code": "{\n  \"provider\": {\n    \"mi-proveedor\": {\n      \"npm\": \"@ai-sdk/openai-compatible\",\n      \"name\": \"Mi Proveedor\",\n      \"options\": {\n        \"baseURL\": \"https://mi-servidor.com/v1\"\n      },\n      \"models\": {\n        \"modelo-custom\": {\n          \"name\": \"Modelo Custom\"\n        }\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "Gestión de Modelos", "bullets": [
                "model: Modelo principal para tareas complejas",
                "small_model: Modelo ligero para títulos y tareas menores",
                "blacklist: Ocultar modelos específicos del selector",
                "whitelist: Mostrar solo modelos específicos"
            ]},
            {"type": "code", "title": "Blacklist y Whitelist", "code": "{\n  \"provider\": {\n    \"anthropic\": {\n      \"blacklist\": [\"claude-opus-4-20250514\"],\n      \"whitelist\": [\"claude-sonnet-4-20250514\"]\n    }\n  }\n}"},
            {"type": "code", "title": "Configuración de Timeouts", "code": "{\n  \"provider\": {\n    \"anthropic\": {\n      \"options\": {\n        \"timeout\": 600000,\n        \"chunkTimeout\": 30000,\n        \"setCacheKey\": true\n      }\n    }\n  }\n}"},
            {"type": "summary", "title": "Resumen", "points": [
                "Proveedores custom para proxies y endpoints",
                "Gestionar modelos visibles con blacklist/whitelist",
                "Optimizar timeouts según necesidades"
            ]}
        ]
    },
    {
        "file": "08-formato-config.pptx",
        "title": "Formato de Configuración",
        "module": 8,
        "duration": "1 hora",
        "objectives": [
            "Entender el formato JSON/JSONC de configuración",
            "Usar el schema para validación",
            "Crear archivos de configuración válidos",
            "Habilitar autocompletado en editores"
        ],
        "slides": [
            {"type": "title", "content": "Clase 8: Formato de Configuración"},
            {"type": "content", "title": "JSON vs JSONC", "bullets": [
                "JSON: Formato estándar, sin comentarios",
                "JSONC: JSON con comentarios (recomendado)",
                "Ambos son soportados por OpenCode",
                "Extensión: opencode.json o opencode.jsonc"
            ]},
            {"type": "code", "title": "Estructura Básica", "code": "{\n  \"$schema\": \"https://opencode.ai/config.json\",\n  \"model\": \"anthropic/claude-sonnet-4-5\",\n  \"small_model\": \"anthropic/claude-haiku-4-5\",\n  \"autoupdate\": true,\n  \"server\": {\n    \"port\": 4096\n  },\n  \"permission\": {\n    \"bash\": \"ask\",\n    \"edit\": \"allow\"\n  }\n}"},
            {"type": "content", "title": "Secciones Principales", "bullets": [
                "model y small_model: Selección de modelos",
                "provider: Configuración de proveedores",
                "permission: Control de permisos",
                "tools: Habilitar/deshabilitar herramientas",
                "mcp: Configuración de MCP servers",
                "agent: Configuración de agentes",
                "server: Configuración del servidor"
            ]},
            {"type": "code", "title": "Autocompletado en VS Code", "code": "// Agregar al settings.json de VS Code\n{\n  \"json.schemas\": [\n    {\n      \"fileMatch\": [\"opencode.json\"],\n      \"url\": \"https://opencode.ai/config.json\"\n    }\n  ]\n}"},
            {"type": "summary", "title": "Resumen", "points": [
                "JSONC recomendado por soporte de comentarios",
                "Schema para validación y autocompletado",
                "Secciones organizadas por funcionalidad"
            ]}
        ]
    },
    {
        "file": "09-ubicaciones-config.pptx",
        "title": "Ubicaciones y Precedencia de Config",
        "module": 9,
        "duration": "1.5 horas",
        "objectives": [
            "Entender las 8 fuentes de configuración",
            "Conocer el orden de precedencia",
            "Configurar en múltiples niveles",
            "Resolver conflictos de configuración"
        ],
        "slides": [
            {"type": "title", "content": "Clase 9: Ubicaciones de Configuración"},
            {"type": "content", "title": "Orden de Precedencia", "bullets": [
                "1. Remote config (.well-known/opencode)",
                "2. Global config (~/.config/opencode/)",
                "3. Custom config (OPENCODE_CONFIG)",
                "4. Project config (opencode.json)",
                "5. .opencode directories",
                "6. Inline config (OPENCODE_CONFIG_CONTENT)",
                "7. Managed settings (por plataforma)",
                "8. macOS managed preferences (MDM)"
            ]},
            {"type": "content", "title": "Configuración Global", "bullets": [
                "Ubicación: ~/.config/opencode/opencode.json",
                "Para preferencias de usuario",
                "Se aplica a todos los proyectos",
                "Override de remote config"
            ]},
            {"type": "content", "title": "Configuración por Proyecto", "bullets": [
                "Ubicación: opencode.json en raíz del proyecto",
                "Para configuración específica del proyecto",
                "Override de global y remote",
                "Seguro para Git (sin credenciales)"
            ]},
            {"type": "content", "title": "Managed Settings", "bullets": [
                "macOS: /Library/Application Support/opencode/",
                "Linux: /etc/opencode/",
                "Windows: %ProgramData%\\opencode",
                "Prioridad máxima - no overrideable",
                "Deploy vía MDM (Jamf, Kandji, FleetDM)"
            ]},
            {"type": "code", "title": "Ejemplo de Merge", "code": "# Global config:\n{\n  \"autoupdate\": true,\n  \"model\": \"anthropic/claude-haiku-4-5\"\n}\n\n# Project config:\n{\n  \"model\": \"anthropic/claude-sonnet-4-5\"\n}\n\n# Resultado final:\n{\n  \"autoupdate\": true,\n  \"model\": \"anthropic/claude-sonnet-4-5\"  // Override\n}"},
            {"type": "summary", "title": "Resumen", "points": [
                "8 fuentes de configuración con precedencia",
                "Project > Global > Remote",
                "Managed settings tienen prioridad máxima",
                "Las configs se mergean, no reemplazan"
            ]}
        ]
    },
    {
        "file": "10-variables-entorno.pptx",
        "title": "Variables de Entorno y Sustitución",
        "module": 10,
        "duration": "1 hora",
        "objectives": [
            "Usar variables de entorno en configuración",
            "Implementar sustitución de variables",
            "Mantener secrets fuera de archivos",
            "Configurar differentes entornos"
        ],
        "slides": [
            {"type": "title", "content": "Clase 10: Variables de Entorno"},
            {"type": "content", "title": "Variables Principales", "bullets": [
                "OPENCODE_CONFIG: Ruta custom de config",
                "OPENCODE_CONFIG_DIR: Directorio custom",
                "OPENCODE_TUI_CONFIG: Ruta de TUI config",
                "OPENCODE_DISABLE_AUTOUPDATE: Desactivar updates",
                "OPENCODE_EXPERIMENTAL: Habilitar experimentales"
            ]},
            {"type": "code", "title": "Sustitución de Variables", "code": "# Environment variables\n{\n  \"model\": \"{env:OPENCODE_MODEL}\",\n  \"provider\": {\n    \"anthropic\": {\n      \"options\": {\n        \"apiKey\": \"{env:ANTHROPIC_API_KEY}\"\n      }\n    }\n  }\n}\n\n# File contents\n{\n  \"instructions\": [\"./custom-instructions.md\"],\n  \"provider\": {\n    \"openai\": {\n      \"options\": {\n        \"apiKey\": \"{file:~/.secrets/openai-key}\"\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "Casos de Uso", "bullets": [
                "API keys en archivos separados (no en config)",
                "Instrucciones grandes en archivos .md",
                "Configuración por ambiente (dev/staging/prod)",
                "Secrets compartidos en equipos"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Variables de entorno para configuración flexible",
                "Sustitución {env:} y {file:} para secrets",
                "Mantener credenciales fuera del repo"
            ]}
        ]
    },
    {
        "file": "11-gestion-remota-enterprise.pptx",
        "title": "Gestión Remota y Enterprise",
        "module": 11,
        "duration": "1 hora",
        "objectives": [
            "Configurar defaults organizacionales",
            "Implementar managed settings con MDM",
            "Entender el sistema de prioridades",
            "Desplegar en organizaciones"
        ],
        "slides": [
            {"type": "title", "content": "Clase 11: Gestión Remota y Enterprise"},
            {"type": "content", "title": "Remote Config", "bullets": [
                "Endpoint: .well-known/opencode",
                "Defaults para toda la organización",
                "Servidores MCP deshabilitados por defecto",
                "Override en config local"
            ]},
            {"type": "code", "title": "Managed Settings - macOS", "code": "<?xml version=\"1.0\" encoding=\"UTF-8\"?>\n<!DOCTYPE plist PUBLIC \"-//Apple//DTD PLIST 1.0//EN\"\n  \"http://www.apple.com/DTDs/PropertyList-1.0.dtd\">\n<plist version=\"1.0\">\n<dict>\n  <key>PayloadContent</key>\n  <array>\n    <dict>\n      <key>PayloadType</key>\n      <string>ai.opencode.managed</string>\n      <key>share</key>\n      <string>disabled</string>\n      <key>permission</key>\n      <dict>\n        <key>bash</key>\n        <string>ask</string>\n      </dict>\n    </dict>\n  </array>\n</dict>\n</plist>"},
            {"type": "content", "title": "Despliegue con MDM", "bullets": [
                "Jamf Pro: Computers > Configuration Profiles",
                "FleetDM: gitops repo bajo mdm.macos_settings",
                "Kandji: Custom Configuration Profile",
                "Verificar: opencode debug config"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Remote config para defaults organizacionales",
                "MDM para enforcement en macOS",
                "Managed settings tienen prioridad máxima"
            ]}
        ]
    },
    {
        "file": "12-navegacion-tui.pptx",
        "title": "Navegación en la Interfaz TUI",
        "module": 12,
        "duration": "1.5 horas",
        "objectives": [
            "Lanzar y navegar OpenCode TUI",
            "Entender los modos Build y Plan",
            "Usar referencias de archivos con @",
            "Manipular imágenes en el terminal"
        ],
        "slides": [
            {"type": "title", "content": "Clase 12: Navegación en la TUI"},
            {"type": "code", "title": "Lanzar OpenCode", "code": "# Navegar al proyecto\ncd /path/to/project\n\n# Lanzar OpenCode\nopencode"},
            {"type": "content", "title": "Elementos de la TUI", "bullets": [
                "Leader key: Ctrl+X (por defecto)",
                "Indicador de modo: Build/Plan en esquina",
                "Área de chat:输入 y respuestas",
                "Barra de estado: modelo, tokens, sesión"
            ]},
            {"type": "content", "title": "Modos de Trabajo", "bullets": [
                "Build mode (default): Acceso completo a tools",
                "Plan mode: Solo lectura, análisis de código",
                "Cambio con Tab o keybind configurado",
                "Plan → Build para implementar"
            ]},
            {"type": "code", "title": "Referencias de Archivos", "code": "# Buscar archivo\n@nombre-archivo\n\n# Referencia específica\n@src/components/Button.tsx\n\n# Líneas específicas\n@src/index.ts#L10-20\n\n# Drag and drop de imágenes"},
            {"type": "summary", "title": "Resumen", "points": [
                "TUI es la interfaz principal",
                "Build para implementar, Plan para analizar",
                "@ para referenciar archivos",
                "Soporte de imágenes con drag and drop"
            ]}
        ]
    },
    {
        "file": "13-slash-commands.pptx",
        "title": "Comandos Slash Completos",
        "module": 13,
        "duration": "1.5 horas",
        "objectives": [
            "Dominar todos los slash commands",
            "Crear flujos de trabajo eficientes",
            "Gestionar sesiones y contexto",
            "Usar comandos de inicialización"
        ],
        "slides": [
            {"type": "title", "content": "Clase 13: Comandos Slash"},
            {"type": "content", "title": "Comandos de Sesión", "bullets": [
                "/new (/clear): Nueva sesión",
                "/sessions (/resume, /continue): Listar sesiones",
                "/compact (/summarize): Compactar contexto",
                "/share: Compartir conversación",
                "/unshare: Revocar enlace",
                "/export: Exportar a Markdown"
            ]},
            {"type": "content", "title": "Comandos de Deshacer", "bullets": [
                "/undo: Deshacer último mensaje",
                "/redo: Rehacer mensaje deshecho",
                "Múltiples undo/redo soportados",
                "Requiere git repository"
            ]},
            {"type": "content", "title": "Comandos de Proveedor", "bullets": [
                "/connect: Agregar/actualizar credenciales",
                "/models: Cambiar modelo activo",
                "Acceso rápido desde el TUI"
            ]},
            {"type": "content", "title": "Comandos de Inicialización", "bullets": [
                "/init: Crear/actualizar AGENTS.md",
                "/init-deep: Inicialización profunda",
                "Analiza estructura del proyecto",
                "Crea contexto para el agente"
            ]},
            {"type": "content", "title": "Comandos de Navegación", "bullets": [
                "/help: Mostrar ayuda",
                "/details: Toggle detalles de ejecución",
                "/thinking: Toggle proceso de razonamiento",
                "/themes: Listar temas",
                "/editor: Abrir editor externo",
                "/exit: Salir de OpenCode"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "20+ slash commands disponibles",
                "/init esencial para cada proyecto",
                "/compact para gestión de contexto",
                "/undo/redo para deshacer cambios"
            ]}
        ]
    },
    {
        "file": "14-atajos-teclado.pptx",
        "title": "Atajos de Teclado",
        "module": 14,
        "duration": "1 hora",
        "objectives": [
            "Memorizar atajos esenciales",
            "Personalizar keybinds",
            "Navegar eficientemente",
            "Aumentar productividad"
        ],
        "slides": [
            {"type": "title", "content": "Clase 14: Atajos de Teclado"},
            {"type": "content", "title": "Leader Key: Ctrl+X", "bullets": [
                "Ctrl+X N: Nueva sesión",
                "Ctrl+X L: Listar sesiones",
                "Ctrl+X C: Compactar",
                "Ctrl+X S: Compartir",
                "Ctrl+X X: Exportar",
                "Ctrl+X U: Deshacer",
                "Ctrl+X R: Rehacer"
            ]},
            {"type": "content", "title": "Más Atajos", "bullets": [
                "Ctrl+X M: Modelos",
                "Ctrl+X E: Editor",
                "Ctrl+X I: Init",
                "Ctrl+X D: Detalles",
                "Ctrl+X Q: Salir",
                "Ctrl+P: Paleta de comandos",
                "Tab: Cambiar modo Build/Plan"
            ]},
            {"type": "code", "title": "Personalizar Keybinds", "code": "# En tui.json\n{\n  \"keybinds\": {\n    \"command_list\": \"ctrl+p\",\n    \"session_new\": \"ctrl+n\",\n    \"session_list\": \"ctrl+l\"\n  }\n}"},
            {"type": "summary", "title": "Resumen", "points": [
                "Ctrl+X como leader key por defecto",
                "Atajos para todas las operaciones comunes",
                "Personalizable en tui.json"
            ]}
        ]
    },
    {
        "file": "15-personalizacion-tui.pptx",
        "title": "Personalización de la TUI",
        "module": 15,
        "duration": "1 hora",
        "objectives": [
            "Configurar temas de color",
            "Personalizar comportamiento del TUI",
            "Habilitar notificaciones",
            "Optimizar la experiencia visual"
        ],
        "slides": [
            {"type": "title", "content": "Clase 15: Personalización de la TUI"},
            {"type": "code", "title": "Estructura de tui.json", "code": "{\n  \"$schema\": \"https://opencode.ai/tui.json\",\n  \"theme\": \"tokyonight\",\n  \"scroll_speed\": 3,\n  \"scroll_acceleration\": {\n    \"enabled\": true\n  },\n  \"diff_style\": \"auto\",\n  \"mouse\": true,\n  \"attention\": {\n    \"enabled\": true,\n    \"notifications\": true,\n    \"sound\": true,\n    \"volume\": 0.4\n  }\n}"},
            {"type": "content", "title": "Opciones de Tema", "bullets": [
                "/themes para listar temas disponibles",
                "theme en tui.json para configurar",
                "Opciones: tokyonight, monokai, solarized, etc.",
                "Auto-detect terminal colors"
            ]},
            {"type": "content", "title": "Scroll y Navegación", "bullets": [
                "scroll_speed: Velocidad de scroll",
                "scroll_acceleration: Aceleración progresiva",
                "mouse: Habilitar soporte de mouse",
                "diff_style: auto, unified, split"
            ]},
            {"type": "content", "title": "Atención y Notificaciones", "bullets": [
                "attention.enabled: Habilitar sistema",
                "notifications: Notificaciones de escritorio",
                "sound: Sonidos de alerta",
                "volume: Volumen (0.0 - 1.0)"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "tui.json para configuración dedicada",
                "Múltiples temas disponibles",
                "Notificaciones para tareas largas",
                "Soporte de mouse opcional"
            ]}
        ]
    },
    {
        "file": "16-tools-principales-parte1.pptx",
        "title": "Herramientas Principales - Parte 1",
        "module": 16,
        "duration": "1.5 horas",
        "objectives": [
            "Dominar las herramientas fundamentales",
            "Usar bash para comandos shell",
            "Editar archivos con precisión",
            "Buscar en el codebase eficientemente"
        ],
        "slides": [
            {"type": "title", "content": "Clase 16: Herramientas Principales (Parte 1)"},
            {"type": "content", "title": "Tool: bash", "bullets": [
                "Ejecutar comandos shell",
                "Instalar dependencias: npm install",
                "Verificar estado: git status",
                "Ejecutar tests: npm test",
                "Cualquier comando del sistema"
            ]},
            {"type": "content", "title": "Tool: edit", "bullets": [
                "Modificación precisa de archivos",
                "Reemplazo de texto exacto",
                "oldString → newString",
                "Edición quirúrgica sin rewrite completo"
            ]},
            {"type": "code", "title": "Ejemplo de edit", "code": "// Cambiar nombre de función\noldString: \"function calculateTotal(items)\"\nnewString: \"function computeTotal(items)\""},
            {"type": "content", "title": "Tool: write", "bullets": [
                "Crear nuevos archivos",
                "Sobrescribir archivos existentes",
                "Controlado por permiso 'edit'",
                "Útil para archivos nuevos"
            ]},
            {"type": "content", "title": "Tool: read", "bullets": [
                "Leer contenido de archivos",
                "Soporte de rangos de líneas",
                "Lectura eficiente de archivos grandes",
                "Análisis de código"
            ]},
            {"type": "content", "title": "Tool: grep", "bullets": [
                "Búsqueda por regex en contenido",
                "Full regex syntax soportado",
                "Filtrado por patrón de archivo",
                "Rápido con ripgrep bajo el capó"
            ]},
            {"type": "content", "title": "Tool: glob", "bullets": [
                "Búsqueda de archivos por patrón",
                "Soporte de glob patterns",
                "Returns archivos ordenados por modificación",
                "Ejemplo: **/*.js, src/**/*.ts"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "6 tools fundamentales para coding",
                "bash para ejecución de comandos",
                "edit/write/read para gestión de archivos",
                "grep/glob para búsqueda"
            ]}
        ]
    },
    {
        "file": "17-tools-principales-parte2.pptx",
        "title": "Herramientas Principales - Parte 2",
        "module": 17,
        "duration": "1.5 horas",
        "objectives": [
            "Usar apply_patch para cambios complejos",
            "Gestionar skills y tareas",
            "Acceder a información web",
            "Entender la herramienta LSP"
        ],
        "slides": [
            {"type": "title", "content": "Clase 17: Herramientas Principales (Parte 2)"},
            {"type": "content", "title": "Tool: apply_patch", "bullets": [
                "Aplicar parches a archivos",
                "Útil para diffs y patches",
                "Cambia archivos múltiples",
                "Controlado por permiso 'edit'"
            ]},
            {"type": "content", "title": "Tool: skill", "bullets": [
                "Cargar archivos SKILL.md",
                "Instrucciones reutilizables",
                "Carga on-demand",
                "Permisos configurables"
            ]},
            {"type": "content", "title": "Tool: todowrite", "bullets": [
                "Gestionar listas de tareas",
                "Crear/actualizar tareas",
                "Seguimiento de progreso",
                "Útil para tareas complejas"
            ]},
            {"type": "content", "title": "Tool: webfetch", "bullets": [
                "Obtener contenido web",
                "Consultar documentación",
                "Investigar recursos online",
                "Soporte markdown y HTML"
            ]},
            {"type": "content", "title": "Tool: websearch", "bullets": [
                "Búsqueda en la web vía Exa AI",
                "Información actualizada",
                "Sin API key requerida",
                "Disponible con OpenCode provider"
            ]},
            {"type": "content", "title": "Tool: question", "bullets": [
                "Preguntar al usuario",
                "Opciones predefinidas",
                "Respuestas custom",
                "Clarificación durante ejecución"
            ]},
            {"type": "content", "title": "Tool: lsp (Experimental)", "bullets": [
                "goToDefinition: Ir a definición",
                "findReferences: Encontrar referencias",
                "hover: Información hover",
                "documentSymbol: Símbolos del documento",
                "Requiere OPENCODE_EXPERIMENTAL=true"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "12+ herramientas disponibles",
                "webfetch/websearch para información",
                "skill para instrucciones reutilizables",
                "lsp para inteligencia de código"
            ]}
        ]
    },
    {
        "file": "18-configuracion-permisos.pptx",
        "title": "Configuración de Permisos",
        "module": 18,
        "duration": "1 hora",
        "objectives": [
            "Entender el sistema de permisos",
            "Configurar permisos por herramienta",
            "Usar patrones glob para múltiples tools",
            "平衡ar seguridad y funcionalidad"
        ],
        "slides": [
            {"type": "title", "content": "Clase 18: Configuración de Permisos"},
            {"type": "content", "title": "Niveles de Permiso", "bullets": [
                "allow: Permitir sin preguntar",
                "ask: Requiere aprobación del usuario",
                "deny: Bloquear completamente",
                "Default: allow (sin restricciones)"
            ]},
            {"type": "code", "title": "Permisos por Tool", "code": "{\n  \"permission\": {\n    \"bash\": \"ask\",\n    \"edit\": \"allow\",\n    \"write\": \"deny\",\n    \"webfetch\": \"allow\",\n    \"skill\": \"ask\"\n  }\n}"},
            {"type": "code", "title": "Patrones Glob", "code": "{\n  \"permission\": {\n    \"mymcp_*\": \"ask\",\n    \"internal-*\": \"deny\",\n    \"experimental-*\": \"ask\"\n  }\n}"},
            {"type": "content", "title": "Recomendaciones", "bullets": [
                "bash: 'ask' en producción",
                "edit/write: 'allow' para desarrollo",
                "MCP tools: 'ask' por seguridad",
                "webfetch/websearch: 'allow'"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "3 niveles: allow, ask, deny",
                "Patrones glob para múltiples tools",
                "Configurar según contexto de seguridad"
            ]}
        ]
    },
    {
        "file": "19-modo-build.pptx",
        "title": "Modo Build",
        "module": 19,
        "duration": "1 hora",
        "objectives": [
            "Entender las capacidades del modo Build",
            "Implementar features directamente",
            "Refactorizar código de forma segura",
            "Usar el modo de manera efectiva"
        ],
        "slides": [
            {"type": "title", "content": "Clase 19: Modo Build"},
            {"type": "content", "title": "¿Qué es el Modo Build?", "bullets": [
                "Modo por defecto de OpenCode",
                "Acceso completo a todas las herramientas",
                "Lectura, escritura y edición de archivos",
                "Ejecución de comandos shell",
                "Búsqueda en el codebase"
            ]},
            {"type": "content", "title": "Casos de Uso", "bullets": [
                "Desarrollo de nuevas features",
                "Refactoring de código existente",
                "Bug fixing con contexto completo",
                "Creación de archivos nuevos",
                "Configuración de proyectos"
            ]},
            {"type": "code", "title": "Ejemplo de Uso", "code": "# Feature directa\n\"We need to add authentication to /settings route.\nMirror how it's handled in /notes route in\n@packages/functions/src/notes.ts and implement\nthe same logic in @packages/functions/src/settings.ts\""},
            {"type": "content", "title": "Mejores Prácticas", "bullets": [
                "Proporcionar contexto suficiente",
                "Referenciar archivos con @",
                "Ser específico en los cambios",
                "Usar /undo si el resultado no es correcto"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Build mode: acceso completo a tools",
                "Ideal para implementación directa",
                "Usar para features y bugs específicos"
            ]}
        ]
    },
    {
        "file": "20-modo-plan.pptx",
        "title": "Modo Plan",
        "module": 20,
        "duration": "1 hora",
        "objectives": [
            "Entender el propósito del modo Plan",
            "Analizar código de forma segura",
            "Proponer cambios sin modificar",
            "Usar Plan para exploración"
        ],
        "slides": [
            {"type": "title", "content": "Clase 20: Modo Plan"},
            {"type": "content", "title": "¿Qué es el Modo Plan?", "bullets": [
                "Modo de solo lectura",
                "Sin capacidad de modificar archivos",
                "Análisis y propuestas de cambios",
                "Ideal para exploración segura"
            ]},
            {"type": "content", "title": "Casos de Uso", "bullets": [
                "Explorar codebase unfamiliar",
                "Planificar features complejas",
                "Revisar sugerencias antes de implementar",
                "Entender arquitectura existente"
            ]},
            {"type": "content", "title": "Flujo de Trabajo", "bullets": [
                "1. Activar Plan mode con Tab",
                "2. Describir lo que se quiere lograr",
                "3. Recibir plan detallado",
                "4. Iterar sobre el plan",
                "5. Cambiar a Build para implementar"
            ]},
            {"type": "code", "title": "Ejemplo de Plan", "code": "# En Plan mode:\n\"When a user deletes a note, we'd like to flag it\nas deleted in the database. Then create a screen\nthat shows all the recently deleted notes. From this\nscreen, the user can undelete a note or permanently\ndelete it.\""},
            {"type": "summary", "title": "Resumen", "points": [
                "Plan mode: solo lectura, análisis",
                "Ideal para exploración y planificación",
                "Cambia a Build cuando estés listo"
            ]}
        ]
    },
    {
        "file": "21-flujo-trabajo-plan-build.pptx",
        "title": "Flujo de Trabajo Plan → Build",
        "module": 21,
        "duration": "1.5 horas",
        "objectives": [
            "Combinar Plan y Build efectivamente",
            "Planificar features complejas",
            "Implementar cambios incrementalmente",
            "Gestionar el ciclo de desarrollo"
        ],
        "slides": [
            {"type": "title", "content": "Clase 21: Flujo de Trabajo Plan → Build"},
            {"type": "content", "title": "El Flujo Completo", "bullets": [
                "1. Activar Plan mode (Tab)",
                "2. Describir feature con detalle",
                "3. Recibir y revisar plan",
                "4. Iterar con feedback",
                "5. Cambiar a Build mode (Tab)",
                "6. Ejecutar: 'Go ahead and make the changes'"
            ]},
            {"type": "content", "title": "Tips para Planificación", "bullets": [
                "Hablar como a un junior developer",
                "Incluir ejemplos y contexto",
                "Referenciar archivos existentes",
                "Ser específico en requisitos"
            ]},
            {"type": "content", "title": "Iteración sobre el Plan", "bullets": [
                "Pedir más detalles si es necesario",
                "Agregar restricciones o preferencias",
                "Comparar alternativas",
                "Asegurar que el plan sea completo"
            ]},
            {"type": "content", "title": "Edición Directa (sin Plan)", "bullets": [
                "Para cambios simples y directos",
                "Referenciar archivos con @",
                "Proporcionar contexto mínimo",
                "Más rápido para tareas triviales"
            ]},
            {"type": "demo", "title": "Demo: Feature Completa", "steps": [
                "Activar Plan mode",
                "Describir feature de autenticación",
                "Recibir plan con pasos",
                "Solicitar ajustes",
                "Cambiar a Build",
                "Implementar cambios",
                "Verificar resultado"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Plan → Build para features complejas",
                "Edición directa para cambios simples",
                "Iterar siempre sobre el plan"
            ]}
        ]
    },
    {
        "file": "22-agentes-primarios-subagentes.pptx",
        "title": "Agentes Primarios y Subagentes",
        "module": 22,
        "duration": "1.5 horas",
        "objectives": [
            "Entender la jerarquía de agentes",
            "Usar agentes build y plan",
            "Invocar subagentes con @mention",
            "Navegar entre sesiones child"
        ],
        "slides": [
            {"type": "title", "content": "Clase 22: Agentes Primarios y Subagentes"},
            {"type": "content", "title": "Agentes Primarios", "bullets": [
                "build (default): Acceso completo a tools",
                "plan: Solo lectura, análisis",
                "Cambio con Tab o keybind",
                "Interacción directa con el usuario"
            ]},
            {"type": "content", "title": "Subagentes", "bullets": [
                "Ejecución en sesiones child",
                "Invocación automática por primarios",
                "Invocación manual con @mention",
                "Profundidad configurable (default: 1)"
            ]},
            {"type": "code", "title": "Invocación de Subagentes", "code": "# Invocación manual\n@general help me search for this function\n\n# El subagente crea una child session\n# Trabaja de forma independiente\n# Retorna resultados al agente primario"},
            {"type": "content", "title": "Navegación entre Sesiones", "bullets": [
                "Leader+Down: Entrar a child session",
                "Right: Ciclar a siguiente child",
                "Left: Ciclar a child anterior",
                "Up: Volver a parent session"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Primarios: build y plan",
                "Subagentes: tareas especializadas",
                "@mention para invocación manual",
                "Navegación con atajos de teclado"
            ]}
        ]
    },
    {
        "file": "23-agentes-custom.pptx",
        "title": "Agentes Custom",
        "module": 23,
        "duration": "1.5 horas",
        "objectives": [
            "Crear agentes personalizados",
            "Configurar herramientas y permisos",
            "Establecer agente por defecto",
            "Diseñar agentes para casos específicos"
        ],
        "slides": [
            {"type": "title", "content": "Clase 23: Agentes Custom"},
            {"type": "code", "title": "Estructura de Agente Custom", "code": "---\ndescription: Performs security audits\nmode: subagent\nmodel: anthropic/claude-sonnet-4-5\ntemperature: 0.1\ntools:\n  write: false\n  edit: false\n  bash: false\n---\n\nYou are a security expert. Focus on identifying\npotential security issues in the code."},
            {"type": "content", "title": "Campos Disponibles", "bullets": [
                "description: Descripción del agente",
                "model: Modelo específico a usar",
                "prompt: Prompt del sistema",
                "temperature: Temperatura (0-1)",
                "tools: Habilitar/deshabilitar tools",
                "permission: Permisos específicos",
                "mode: primary o subagent"
            ]},
            {"type": "code", "title": "Config en opencode.json", "code": "{\n  \"agent\": {\n    \"code-reviewer\": {\n      \"description\": \"Reviews code for best practices\",\n      \"model\": \"anthropic/claude-sonnet-4-5\",\n      \"prompt\": \"You are a code reviewer.\",\n      \"tools\": {\n        \"write\": false,\n        \"edit\": false\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "Agente por Defecto", "bullets": [
                "Configurar: default_agent en opencode.json",
                "Aplica a: TUI, CLI, desktop, GitHub Action",
                "Fallback a build si no existe",
                "Debe ser agente primario"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Agentes custom para tareas específicas",
                "YAML frontmatter o JSON en config",
                "Control total de tools y permisos",
                "default_agent para cambio global"
            ]}
        ]
    },
    {
        "file": "24-agent-skills.pptx",
        "title": "Agent Skills (SKILL.md)",
        "module": 24,
        "duration": "1.5 horas",
        "objectives": [
            "Crear skills reutilizables",
            "Entender la estructura SKILL.md",
            "Configurar permisos por skill",
            "Usar skills en flujos de trabajo"
        ],
        "slides": [
            {"type": "title", "content": "Clase 24: Agent Skills"},
            {"type": "content", "title": "¿Qué son las Skills?", "bullets": [
                "Instrucciones reutilizables",
                "Descubiertas por agentes",
                "Carga on-demand via tool skill",
                "Archivos SKILL.md"
            ]},
            {"type": "code", "title": "Ubicaciones", "code": "# Proyecto\n.opencode/skills/<name>/SKILL.md\n\n# Global\n~/.config/opencode/skills/<name>/SKILL.md\n\n# Claude-compatible\n.claude/skills/<name>/SKILL.md\n\n# Agent-compatible\n.agents/skills/<name>/SKILL.md"},
            {"type": "code", "title": "Estructura SKILL.md", "code": "---\nname: git-release\ndescription: Create consistent releases\nlicense: MIT\ncompatibility: opencode\nmetadata:\n  audience: maintainers\n---\n\n## What I do\n- Draft release notes from merged PRs\n- Propose semantic version bumps\n\n## When to use me\nUse this when preparing tagged releases."},
            {"type": "code", "title": "Permisos de Skills", "code": "{\n  \"permission\": {\n    \"skill\": {\n      \"*\": \"allow\",\n      \"internal-*\": \"deny\",\n      \"experimental-*\": \"ask\"\n    }\n  }\n}"},
            {"type": "summary", "title": "Resumen", "points": [
                "Skills: instrucciones reutilizables",
                "SKILL.md con YAML frontmatter",
                "Permisos granulares por skill",
                "Carga on-demand eficiente"
            ]}
        ]
    },
    {
        "file": "25-agents-md-contexto.pptx",
        "title": "AGENTS.md y Contexto del Proyecto",
        "module": 25,
        "duration": "1 hora",
        "objectives": [
            "Crear AGENTS.md con /init",
            "Entender su importancia para agentes",
            "Editar manualmente el contexto",
            "Mantener AGENTS.md actualizado"
        ],
        "slides": [
            {"type": "title", "content": "Clase 25: AGENTS.md"},
            {"type": "content", "title": "¿Qué es AGENTS.md?", "bullets": [
                "Archivo de contexto del proyecto",
                "Creado con /init",
                "Analiza estructura y patrones",
                "Persiste entre sesiones"
            ]},
            {"type": "code", "title": "Crear con /init", "code": "# En OpenCode TUI\n/init\n\n# Esto:\n# 1. Analiza la estructura del proyecto\n# 2. Identifica patrones y convenciones\n# 3. Crea AGENTS.md en la raíz\n# 4. Se recomienda commitear a Git"},
            {"type": "content", "title": "Contenido Típico", "bullets": [
                "Estructura del proyecto",
                "Convenciones de naming",
                "Patrones de arquitectura",
                "Comandos disponibles",
                "Guías de coding"
            ]},
            {"type": "content", "title": "Instrucciones Adicionales", "bullets": [
                "Campo instructions en config",
                "Array de paths y globs",
                "Se carga junto con AGENTS.md",
                "Ejemplo: CONTRIBUTING.md, guidelines"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "AGENTS.md: contexto persistente del proyecto",
                "/init para crear automáticamente",
                "Commitear a Git para compartir",
                "Complementar con instructions"
            ]}
        ]
    },
    {
        "file": "26-mcp-servers.pptx",
        "title": "MCP Servers",
        "module": 26,
        "duration": "1.5 horas",
        "objectives": [
            "Entender el Model Context Protocol",
            "Configurar MCP local y remoto",
            "Implementar OAuth para MCP",
            "Gestionar MCP por agente"
        ],
        "slides": [
            {"type": "title", "content": "Clase 26: MCP Servers"},
            {"type": "content", "title": "¿Qué es MCP?", "bullets": [
                "Model Context Protocol",
                "Protocolo para herramientas externas",
                "Soporte local y remoto",
                "Tools disponibles junto a tools nativas"
            ]},
            {"type": "code", "title": "MCP Local", "code": "{\n  \"mcp\": {\n    \"my-local-mcp\": {\n      \"type\": \"local\",\n      \"command\": [\"npx\", \"-y\", \"my-mcp-command\"],\n      \"enabled\": true,\n      \"environment\": {\n        \"MY_ENV_VAR\": \"value\"\n      }\n    }\n  }\n}"},
            {"type": "code", "title": "MCP Remoto", "code": "{\n  \"mcp\": {\n    \"my-remote-mcp\": {\n      \"type\": \"remote\",\n      \"url\": \"https://my-mcp-server.com\",\n      \"enabled\": true,\n      \"headers\": {\n        \"Authorization\": \"Bearer TOKEN\"\n      }\n    }\n  }\n}"},
            {"type": "content", "title": "OAuth para MCP Remoto", "bullets": [
                "Detección automática de 401",
                "Dynamic Client Registration",
                "Tokens en ~/.local/share/opencode/",
                "Comandos: mcp auth, list, logout, debug"
            ]},
            {"type": "content", "title": "Ejemplos Populares", "bullets": [
                "Sentry: Errores y issues",
                "Context7: Búsqueda de docs",
                "Grep by Vercel: Código en GitHub"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "MCP: protocolo para tools externas",
                "Local con command, remoto con URL",
                "OAuth automático para remoto",
                "Gestión por agente disponible"
            ]}
        ]
    },
    {
        "file": "27-custom-tools-plugins.pptx",
        "title": "Custom Tools y Plugins",
        "module": 27,
        "duration": "1 hora",
        "objectives": [
            "Crear herramientas personalizadas",
            "Usar el sistema de plugins",
            "Integrar herramientas externas",
            "Extender funcionalidad de OpenCode"
        ],
        "slides": [
            {"type": "title", "content": "Clase 27: Custom Tools y Plugins"},
            {"type": "content", "title": "Custom Tools", "bullets": [
                "Funciones definidas por el usuario",
                "Ejecución de código arbitrario",
                "Definidos en config file",
                "Disponibles para el LLM"
            ]},
            {"type": "content", "title": "Sistema de Plugins", "bullets": [
                "Ubicación: .opencode/plugins/",
                "Global: ~/.config/opencode/plugins/",
                "Carga desde npm",
                "Extienden con tools, hooks, integraciones"
            ]},
            {"type": "code", "title": "Configuración de Plugins", "code": "{\n  \"plugin\": [\n    \"opencode-helicone-session\",\n    \"@my-org/custom-plugin\"\n  ]\n}"},
            {"type": "content", "title": "Ejemplos de Plugins", "bullets": [
                "opencode-helicone-session: Tracking de sesiones",
                "opencode-gitlab-plugin: Herramientas GitLab",
                "Plugins community en npm"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Custom tools para funcionalidad específica",
                "Plugins para extensiones completas",
                "Carga desde npm o directorios locales"
            ]}
        ]
    },
    {
        "file": "28-integracion-ide.pptx",
        "title": "Integración IDE",
        "module": 28,
        "duration": "1 hora",
        "objectives": [
            "Instalar extensión en VS Code/Cursor",
            "Usar atajos de la extensión",
            "Compartir contexto con el IDE",
            "Configurar editor externo"
        ],
        "slides": [
            {"type": "title", "content": "Clase 28: Integración IDE"},
            {"type": "content", "title": "IDEs Soportados", "bullets": [
                "VS Code",
                "Cursor",
                "Windsurf",
                "VSCodium",
                "Cualquiera con terminal integrado"
            ]},
            {"type": "content", "title": "Instalación", "bullets": [
                "Automática: ejecutar opencode en terminal",
                "Manual: buscar OpenCode en Marketplace",
                "Requiere CLI del IDE en PATH"
            ]},
            {"type": "content", "title": "Atajos de la Extensión", "bullets": [
                "Cmd/Ctrl+Esc: Abrir/focus OpenCode",
                "Cmd/Ctrl+Shift+Esc: Nueva sesión",
                "Cmd/Ctrl+Option+K: Insertar referencia"
            ]},
            {"type": "content", "title": "Context Awareness", "bullets": [
                "Comparte selección actual",
                "Comparte tab abierto",
                "Referencias de archivo @File#L37-42",
                "Drag and drop de imágenes"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Extensión para VS Code y forks",
                "Instalación automática o manual",
                "Atajos paraproductividad",
                "Context awareness integrado"
            ]}
        ]
    },
    {
        "file": "29-integracion-github.pptx",
        "title": "Integración GitHub",
        "module": 29,
        "duration": "1.5 horas",
        "objectives": [
            "Configurar GitHub Actions con OpenCode",
            "Automatizar review de PRs",
            "Triage automático de issues",
            "Usar /opencode en comentarios"
        ],
        "slides": [
            {"type": "title", "content": "Clase 29: Integración GitHub"},
            {"type": "code", "title": "Instalación", "code": "# Instalar GitHub app y workflow\nopencode github install\n\n# O manualmente:\# 1. Instalar app: github.com/apps/opencode-agent\n# 2. Crear workflow: .github/workflows/opencode.yml"},
            {"type": "code", "title": "Workflow Básico", "code": "name: opencode\non:\n  issue_comment:\n    types: [created]\n  pull_request_review_comment:\n    types: [created]\njobs:\n  opencode:\n    if: |\n      contains(github.event.comment.body, '/oc') ||\n      contains(github.event.comment.body, '/opencode')\n    runs-on: ubuntu-latest\n    steps:\n      - uses: actions/checkout@v6\n      - uses: anomalyco/opencode/github@latest\n        env:\n          ANTHROPIC_API_KEY: ${{ secrets.ANTHROPIC_API_KEY }}\n        with:\n          model: anthropic/claude-sonnet-4-5"},
            {"type": "content", "title": "Eventos Soportados", "bullets": [
                "issue_comment: Comentarios en issues/PRs",
                "pull_request_review_comment: Líneas específicas",
                "issues: Issues abiertos/editados",
                "pull_request: PRs abiertos/actualizados",
                "schedule: Tareas cron",
                "workflow_dispatch: Trigger manual"
            ]},
            {"type": "content", "title": "Uso en GitHub", "bullets": [
                "/opencode explain this issue",
                "/opencode fix this",
                "/oc delete attachment",
                "Comentarios en líneas de código"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "GitHub Actions para automatización",
                "Múltiples eventos soportados",
                "/opencode para interacción",
                "Runs en GitHub runners (seguro)"
            ]}
        ]
    },
    {
        "file": "30-integracion-gitlab.pptx",
        "title": "Integración GitLab",
        "module": 30,
        "duration": "1 hora",
        "objectives": [
            "Configurar GitLab Duo Agent Platform",
            "Usar OAuth o Personal Access Token",
            "Integrar con self-hosted GitLab",
            "Usar el plugin de GitLab"
        ],
        "slides": [
            {"type": "title", "content": "Clase 30: Integración GitLab"},
            {"type": "content", "title": "GitLab Duo Agent Platform", "bullets": [
                "Requiere Premium o Ultimate",
                "OAuth o Personal Access Token",
                "Modelos Claude-based",
                "Soporte para GitLab.com y Self-Managed"
            ]},
            {"type": "code", "title": "Configuración", "code": "# Conectar con GitLab\n/connect\n# Seleccionar GitLab\n# Elegir OAuth o Personal Access Token\n\n# Modelos disponibles:\n# - duo-chat-haiku-4-5 (rápido)\n# - duo-chat-sonnet-4-5 (balanceado)\n# - duo-chat-opus-4-5 (capaz)"},
            {"type": "code", "title": "Self-Hosted GitLab", "code": "export GITLAB_INSTANCE_URL=https://gitlab.company.com\nexport GITLAB_TOKEN=glpat-...\nexport GITLAB_AI_GATEWAY_URL=https://ai-gateway.company.com"},
            {"type": "code", "title": "Plugin GitLab", "code": "{\n  \"plugin\": [\"opencode-gitlab-plugin\"]\n}\n\n# Herramientas incluidas:\n# - MR reviews\n# - Issue tracking\n# - Pipeline monitoring\n# - CI/CD integration"},
            {"type": "summary", "title": "Resumen", "points": [
                "GitLab Duo para agentes AI",
                "Premium/Ultimate requerido",
                "Self-hosted soportado",
                "Plugin para herramientas completas"
            ]}
        ]
    },
    {
        "file": "31-troubleshooting-recursos.pptx",
        "title": "Troubleshooting y Recursos",
        "module": 31,
        "duration": "1 hora",
        "objectives": [
            "Resolver problemas comunes",
            "Usar comandos de debug",
            "Encontrar recursos de ayuda",
            "Mantener OpenCode actualizado"
        ],
        "slides": [
            {"type": "title", "content": "Clase 31: Troubleshooting y Recursos"},
            {"type": "content", "title": "Problemas Comunes", "bullets": [
                "Windows: Usar WSL para mejor experiencia",
                "MCP no carga: Verificar permisos",
                "Skill no aparece: Nombre en MAYÚSCULAS",
                "LSP no funciona: Habilitar experimental",
                "Imágenes: Verificar límites de tamaño"
            ]},
            {"type": "code", "title": "Comandos de Debug", "code": "# Ver configuración resuelta\nopencode debug config\n\n# Listar MCP servers y auth\nopencode mcp auth list\n\n# Debug de MCP específico\nopencode mcp debug <server-name>\n\n# Verificar versión\nopencode --version"},
            {"type": "content", "title": "Recursos Oficiales", "bullets": [
                "Docs: opencode.ai/docs",
                "GitHub: github.com/anomalyco/opencode",
                "Discord: opencode.ai/discord",
                "Changelog: opencode.ai/changelog"
            ]},
            {"type": "content", "title": "Comunidad", "bullets": [
                "Discord para soporte",
                "GitHub Issues para bugs",
                "X (Twitter): @opencode",
                "Contribuciones bienvenidas"
            ]},
            {"type": "content", "title": "Actualizaciones", "bullets": [
                "autoupdate: true en config",
                "Verificar: opencode --version",
                "Changelog para ver cambios",
                "Probar en environment de staging"
            ]},
            {"type": "summary", "title": "Resumen", "points": [
                "Problemas comunes tienen soluciones simples",
                "Comandos de debug para diagnóstico",
                "Comunidad activa para soporte",
                "Mantener actualizado es importante"
            ]}
        ]
    }
]

def create_presentation(pres_data, output_dir):
    """Crea una presentación PowerPoint"""
    prs = Presentation()
    
    # Configurar tamaño de diapositiva (widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in pres_data["slides"]:
        slide_type = slide_data["type"]
        
        if slide_type == "title":
            # Diapositiva de título
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            add_title_slide(slide, slide_data["content"], pres_data["module"], pres_data["duration"])
            
        elif slide_type == "content":
            # Diapositiva de contenido con viñetas
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_content_slide(slide, slide_data["title"], slide_data["bullets"])
            
        elif slide_type == "code":
            # Diapositiva de código
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_code_slide(slide, slide_data["title"], slide_data["code"])
            
        elif slide_type == "comparison":
            # Diapositiva de comparación/tabla
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_comparison_slide(slide, slide_data["title"], slide_data["headers"], slide_data["rows"])
            
        elif slide_type == "demo":
            # Diapositiva de demo
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_demo_slide(slide, slide_data["title"], slide_data["steps"])
            
        elif slide_type == "summary":
            # Diapositiva de resumen
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_summary_slide(slide, slide_data["title"], slide_data["points"])
    
    # Guardar
    output_path = os.path.join(output_dir, pres_data["file"])
    prs.save(output_path)
    print(f"Creado: {pres_data['file']}")

def add_title_slide(slide, title_text, module, duration):
    """Agrega diapositiva de título"""
    # Fondo oscuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo con módulo y duración
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Módulo {module} | Duración: {duration}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(slide, title_text, bullets):
    """Agrega diapositiva de contenido"""
    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Línea separadora
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['secondary']
    line.line.fill.background()
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['primary']
        p.space_after = Pt(12)

def add_code_slide(slide, title_text, code_text):
    """Agrega diapositiva de código"""
    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Fondo del código
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = COLORS['code_bg']
    code_bg.line.fill.background()
    
    # Código
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.1))
    tf = code_box.text_frame
    tf.word_wrap = True
    
    lines = code_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.name = "Consolas"
        p.font.color.rgb = COLORS['green']

def add_comparison_slide(slide, title_text, headers, rows):
    """Agrega diapositiva de comparación/tabla"""
    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # Tabla
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5))
    table = table_shape.table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['primary']
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = COLORS['white']
        p.font.bold = True
        p.font.size = Pt(16)
    
    # Rows
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = cell_text
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light_gray']
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)

def add_demo_slide(slide, title_text, steps):
    """Agrega diapositiva de demo"""
    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Demo: {title_text}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Pasos
    for i, step in enumerate(steps):
        step_box = slide.shapes.add_textbox(Inches(1), Inches(1.5 + i * 0.8), Inches(11), Inches(0.7))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {step}"
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['accent']

def add_summary_slide(slide, title_text, points):
    """Agrega diapositiva de resumen"""
    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Puntos
    for i, point in enumerate(points):
        point_box = slide.shapes.add_textbox(Inches(1), Inches(1.5 + i * 0.9), Inches(11), Inches(0.8))
        tf = point_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['green']

def main():
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_dir = os.path.join(base_dir, "presentaciones")
    os.makedirs(output_dir, exist_ok=True)
    
    print("Generando 31 presentaciones PowerPoint...")
    print("=" * 50)
    
    for i, pres_data in enumerate(PRESENTATIONS, 1):
        try:
            create_presentation(pres_data, output_dir)
        except Exception as e:
            print(f"Error en {pres_data['file']}: {e}")
    
    print("=" * 50)
    print("¡Todas las presentaciones generadas!")
    print(f"Ubicación: {output_dir}")

if __name__ == "__main__":
    main()
