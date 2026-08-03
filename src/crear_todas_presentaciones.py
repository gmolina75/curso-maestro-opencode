#!/usr/bin/env python3
"""
Script para crear todas las presentaciones del curso OpenCode
usando el estilo de la plantilla (Roboto Slab + colores azul/navy)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import re

# ==================== COLORES Y FUENTES ====================
COLOR_TITLE = RGBColor(0x32, 0x57, 0xB8)    # #3257B8 - Azul
COLOR_BODY = RGBColor(0x15, 0x21, 0x3F)     # #15213F - Navy oscuro
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA) # Fondo claro
COLOR_ACCENT = RGBColor(0x32, 0x57, 0xB8)   # Mismo azul para acentos
COLOR_TABLE_HEADER = RGBColor(0x32, 0x57, 0xB8)
COLOR_TABLE_ALT = RGBColor(0xE8, 0xEB, 0xF2)

FONT_TITLE = 'Roboto Slab'
FONT_BODY = 'Roboto'

# ==================== FUNCIONES AUXILIARES ====================

def parse_markdown_table(lines):
    """Parsear una tabla markdown."""
    rows = []
    for line in lines:
        if '|' in line and not line.strip().startswith('|---'):
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells:
                rows.append(cells)
    return rows

def extract_sections_from_markdown(content):
    """Extraer secciones del markdown."""
    sections = []
    current_section = None
    current_content = []
    in_code_block = False
    
    lines = content.split('\n')
    
    for line in lines:
        # Detectar bloques de codigo
        if line.strip().startswith('```'):
            in_code_block = not in_code_block
            continue
        
        if in_code_block:
            current_content.append(line)
            continue
        
        # Detectar titulares (## )
        if line.startswith('## '):
            if current_section:
                sections.append({
                    'title': current_section,
                    'content': current_content
                })
            current_section = line.replace('## ', '').strip()
            current_content = []
        elif line.strip():
            current_content.append(line)
    
    # Agregar ultima seccion
    if current_section:
        sections.append({
            'title': current_section,
            'content': current_content
        })
    
    return sections

def extract_tables_from_content(content):
    """Extraer tablas del contenido."""
    tables = []
    lines = content.split('\n')
    current_table = []
    in_table = False
    
    for line in lines:
        if '|' in line:
            in_table = True
            current_table.append(line)
        else:
            if in_table and current_table:
                tables.append(current_table)
                current_table = []
            in_table = False
    
    if current_table:
        tables.append(current_table)
    
    return tables

def extract_bullets_from_content(content):
    """Extraer viñetas del contenido."""
    bullets = []
    for line in content:
        line = line.strip()
        if line.startswith('- ') or line.startswith('* '):
            bullet = line[2:].strip()
            # Limpiar markdown bold
            bullet = re.sub(r'\*\*(.*?)\*\*', r'\1', bullet)
            bullets.append(bullet)
    return bullets

def extract_title_from_markdown(content):
    """Extraer el titulo principal del markdown."""
    for line in content.split('\n'):
        if line.startswith('# ') and not line.startswith('## '):
            return line.replace('# ', '').strip()
    return None

# ==================== FUNCIONES DE CREACION DE DIAPOSITIVAS ====================

def create_title_slide(prs, title, subtitle, author="Curso OpenCode"):
    """Crear diapositiva de titulo."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fondo blanco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(5.5), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(38)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(5.5), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BODY
        p.alignment = PP_ALIGN.LEFT
    
    # Autor
    author_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(5.5), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    p.text = author
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_BODY
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def create_content_slide(prs, title, content_items):
    """Crear diapositiva de contenido."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.5), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(10)
    
    return slide

def create_table_slide(prs, title, table_data):
    """Crear diapositiva con tabla."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    if table_data and len(table_data) > 0:
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows > 0 and cols > 0:
            # Calcular dimensiones
            table_width = min(8.8, cols * 2.2)
            table_height = min(4.8, rows * 0.7)
            
            left = Inches((10 - table_width) / 2)
            top = Inches(1.3)
            
            table_shape = slide.shapes.add_table(
                rows, cols, left, top,
                Inches(table_width), Inches(table_height)
            )
            table = table_shape.table
            
            # Llenar tabla
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_text in enumerate(row_data):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        # Limpiar markdown
                        clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', cell_text)
                        cell.text = clean_text
                        
                        for para in cell.text_frame.paragraphs:
                            para.font.name = FONT_BODY
                            para.font.size = Pt(11)
                            para.font.color.rgb = COLOR_BODY
                        
                        # Header row
                        if row_idx == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
                            for para in cell.text_frame.paragraphs:
                                para.font.color.rgb = COLOR_WHITE
                                para.font.bold = True
                        # Alternating rows
                        elif row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLOR_TABLE_ALT
    
    return slide

def create_bullets_slide(prs, title, bullets):
    """Crear diapositiva con viñetas."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_TITLE
    
    # Viñetas
    bullets_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.5), Inches(4.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.name = FONT_BODY
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(12)
    
    return slide

def create_section_slide(prs, title):
    """Crear diapositiva de seccion."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Titulo centrado
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_closing_slide(prs, title="Gracias", subtitle="Preguntas?"):
    """Crear diapositiva de cierre."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG
    
    # Titulo
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_TITLE
    p.font.size = Pt(44)
    p.font.color.rgb = COLOR_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitulo
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = FONT_BODY
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_BODY
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== PROCESAMIENTO PRINCIPAL ====================

def process_paper(paper_path, output_path):
    """Procesar un paper y crear su presentacion."""
    # Leer el paper
    with open(paper_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Extraer informacion
    main_title = extract_title_from_markdown(content)
    sections = extract_sections_from_markdown(content)
    
    if not main_title:
        main_title = os.path.splitext(os.path.basename(paper_path))[0]
    
    # Crear presentacion
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Extraer metadata del frontmatter
    subtitle = ""
    duration = ""
    for line in content.split('\n')[:10]:
        if 'module:' in line:
            module = line.split(':')[1].strip()
            subtitle = f"Modulo {module}"
        if 'duration:' in line:
            duration = line.split(':')[1].strip().strip('"')
            if subtitle:
                subtitle += f" | Duracion: {duration}"
            else:
                subtitle = f"Duracion: {duration}"
    
    # Diapositiva 1: Titulo
    create_title_slide(prs, main_title, subtitle)
    
    # Procesar cada seccion
    for section in sections:
        section_title = section['title']
        section_content = section['content']
        
        # Verificar si hay tablas
        tables = extract_tables_from_content('\n'.join(section_content))
        
        if tables:
            # Crear diapositiva con tabla
            for table_lines in tables:
                table_data = parse_markdown_table(table_lines)
                if table_data:
                    create_table_slide(prs, section_title, table_data)
        else:
            # Extraer viñetas
            bullets = extract_bullets_from_content(section_content)
            
            if bullets:
                # Crear diapositiva con viñetas
                create_bullets_slide(prs, section_title, bullets[:8])  # Max 8 bullets
            else:
                # Crear diapositiva de contenido regular
                # Filtrar lineas vacias y comentarios
                content_items = []
                for line in section_content:
                    line = line.strip()
                    if line and not line.startswith('#') and not line.startswith('```'):
                        # Limpiar markdown
                        clean_line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
                        if clean_line:
                            content_items.append(clean_line)
                
                if content_items:
                    create_content_slide(prs, section_title, content_items[:10])  # Max 10 items
    
    # Diapositiva de cierre
    create_closing_slide(prs)
    
    # Guardar
    prs.save(output_path)
    return output_path

def main():
    # Directorios relativos al script (src/generar_presentaciones.py)
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    papers_dir = os.path.join(base_dir, 'papers')
    output_dir = os.path.join(base_dir, 'presentaciones')
    
    # Crear directorio de salida si no existe
    os.makedirs(output_dir, exist_ok=True)
    
    # Obtener todos los papers
    papers = sorted([f for f in os.listdir(papers_dir) if f.endswith('.md')])
    
    print(f"Procesando {len(papers)} papers...")
    print()
    
    for paper_file in papers:
        paper_path = os.path.join(papers_dir, paper_file)
        output_file = os.path.join(output_dir, paper_file.replace('.md', '.pptx'))
        
        try:
            process_paper(paper_path, output_file)
            print(f"OK: {paper_file}")
        except Exception as e:
            print(f"ERROR: {paper_file}: {e}")
    
    print()
    print("Proceso completado!")

if __name__ == "__main__":
    main()
